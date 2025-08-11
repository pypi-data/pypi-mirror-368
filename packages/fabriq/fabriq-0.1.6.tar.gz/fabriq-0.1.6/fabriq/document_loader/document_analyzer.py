import traceback
from typing import Dict, Optional, Any
from abc import ABC, abstractmethod
from enum import Enum
import logging
import fitz
from docx import Document as DocxDocument
import openpyxl
from pptx import Presentation


class DocumentContent:
    """Represents analyzed document content characteristics"""

    has_text: bool = False
    has_tables: bool = False
    has_images: bool = False
    has_charts: bool = False
    has_diagrams: bool = False
    has_forms: bool = False
    has_handwriting: bool = False
    has_complex_layout: bool = False
    text_density: float = 0.0
    image_count: int = 0
    table_count: int = 0
    page_count: int = 0
    is_scanned: bool = False
    language: Optional[str] = None


class DocumentType(Enum):
    """Document type enumeration"""

    PDF = "pdf"
    DOCX = "docx"
    DOC = "doc"
    XLSX = "xlsx"
    XLS = "xls"
    PPTX = "pptx"
    PPT = "ppt"
    TXT = "txt"
    RTF = "rtf"
    HTML = "html"
    CSV = "csv"
    JSON = "json"
    XML = "xml"
    # Media files
    MP4 = "mp4"
    AVI = "avi"
    MKV = "mkv"
    MOV = "mov"
    MP3 = "mp3"
    WAV = "wav"
    M4A = "m4a"
    # Images
    PNG = "png"
    JPG = "jpg"
    JPEG = "jpeg"
    GIF = "gif"
    BMP = "bmp"
    TIFF = "tiff"
    WEBP = "webp"


class LoaderType(Enum):
    """Available document loader types"""

    PYMUPDF4LLM = "pymupdf4llm"
    # DOCLING = "docling"
    UNSTRUCTURED = "unstructured"
    MARKITDOWN = "markitdown"
    WHISPER = "whisper"
    PANDAS = "pandas"
    OPENPYXL = "openpyxl"
    PYTHON_DOCX = "python_docx"
    PPTX = "pptx"
    CUSTOM_OCR = "custom_ocr"
    MULTIMODAL_LLM = "multimodal_llm"
    SIMPLE_TEXT = "simple_text"


class BaseLoader(ABC):
    """Base class for all document loaders"""

    @abstractmethod
    def load(self, file_path: str) -> Dict[str, Any]:
        pass

    @abstractmethod
    def can_handle(self, doc_type: DocumentType, content: DocumentContent) -> bool:
        pass


class ContentAnalyzer:
    """Analyzes document content to determine optimal loader"""

    def __init__(self):
        self.logger = logging.getLogger(__name__)

    def analyze_pdf(self, file_path: str) -> DocumentContent:
        """Analyze PDF content characteristics"""
        content = DocumentContent()

        try:
            doc = fitz.open(file_path)
            content.page_count = len(doc)

            total_text_chars = 0
            total_area = 0

            for page_num in range(len(doc)):
                page = doc[page_num]

                # Text analysis
                text = page.get_text()
                total_text_chars += len(text.strip())

                # Get page dimensions
                rect = page.rect
                page_area = rect.width * rect.height
                total_area += page_area

                # Check for tables
                table_finder = page.find_tables()
                tables = table_finder.tables
                if tables:
                    content.has_tables = True
                    content.table_count += len(tables)

                # Check for images
                images = page.get_images()
                if images:
                    content.has_images = True
                    content.image_count += len(images)

                # Check for drawings/diagrams
                drawings = page.get_drawings()
                if drawings:
                    content.has_diagrams = True

                # Check for forms
                widgets = page.widgets()
                if widgets:
                    content.has_forms = True

                # Check for complex layout (multiple columns, etc.)
                blocks = page.get_text("dict")["blocks"]
                if self._has_complex_layout(blocks):
                    content.has_complex_layout = True

            # Calculate text density
            if total_area > 0:
                content.text_density = total_text_chars / total_area

            # Determine if document is scanned
            content.is_scanned = self._is_scanned_pdf(doc)

            # Check if has meaningful text
            content.has_text = total_text_chars > 100

            doc.close()

        except Exception as e:
            self.logger.error(f"Error analyzing PDF {file_path}: {e}\n{traceback.format_exc()}")

        return content

    def analyze_docx(self, file_path: str) -> DocumentContent:
        """Analyze DOCX content characteristics"""
        content = DocumentContent()

        try:
            doc = DocxDocument(file_path)

            # Count text
            text_content = ""
            for paragraph in doc.paragraphs:
                text_content += paragraph.text

            content.has_text = len(text_content.strip()) > 0

            # Count tables
            content.table_count = len(doc.tables)
            content.has_tables = content.table_count > 0

            # Count images
            from docx.oxml.ns import qn

            image_count = 0
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    image_count += 1

            content.image_count = image_count
            content.has_images = image_count > 0

        except Exception as e:
            self.logger.error(f"Error analyzing DOCX {file_path}: {e}")

        return content

    def analyze_excel(self, file_path: str) -> DocumentContent:
        """Analyze Excel content characteristics"""
        content = DocumentContent()

        try:
            wb = openpyxl.load_workbook(file_path)

            # Check for charts
            for sheet in wb.worksheets:
                if sheet._charts:
                    content.has_charts = True
                    break

            # Check for images
            for sheet in wb.worksheets:
                if sheet._images:
                    content.has_images = True
                    content.image_count += len(sheet._images)

            content.has_tables = True  # Excel is inherently tabular

        except Exception as e:
            self.logger.error(f"Error analyzing Excel {file_path}: {e}")

        return content

    def analyze_pptx(self, file_path: str) -> DocumentContent:
        """Analyze PowerPoint content characteristics"""
        content = DocumentContent()

        try:
            prs = Presentation(file_path)

            text_content = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content += shape.text

            content.has_text = len(text_content.strip()) > 0

            # Count images and charts
            image_count = 0
            chart_count = 0

            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.shape_type == 13:  # Picture
                        image_count += 1
                    elif shape.shape_type == 3:  # Chart
                        chart_count += 1

            content.has_images = image_count > 0
            content.image_count = image_count
            content.has_charts = chart_count > 0

        except Exception as e:
            self.logger.error(f"Error analyzing PPTX {file_path}: {e}")

        return content

    def _has_complex_layout(self, blocks) -> bool:
        """Check if document has complex layout"""
        if len(blocks) < 2:
            return False

        # Check for multiple columns
        x_positions = []
        for block in blocks:
            if block.get("type") == 0:  # Text block
                x_positions.append(block["bbox"][0])

        if len(set(x_positions)) > 2:  # More than 2 distinct x positions
            return True

        return False

    def _is_scanned_pdf(self, doc) -> bool:
        """Determine if PDF is scanned"""
        text_to_image_ratio = 0
        total_pages = len(doc)

        for page_num in range(min(3, total_pages)):  # Check first 3 pages
            page = doc[page_num]
            text_length = len(page.get_text().strip())
            image_count = len(page.get_images())

            if image_count > 0 and text_length < 100:
                text_to_image_ratio += 1

        return text_to_image_ratio / min(3, total_pages) > 0.5
