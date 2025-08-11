import pymupdf4llm
import markitdown
from unstructured.partition.auto import partition
import whisper
import openpyxl
import pandas as pd
from pptx import Presentation
from fabriq.document_loader.document_analyzer import BaseLoader, DocumentType, DocumentContent
from typing import Dict, Any
from tabulate import tabulate
from io import StringIO
from langchain.schema import Document
import os


class PyMuPDF4LLMLoader(BaseLoader):
    def load(self, file_path: str) -> Dict[str, Any]:
        content = pymupdf4llm.to_markdown(
            file_path,
            page_chunks=True,
            embed_images=True,
            image_size_limit=0,
            force_text=False,
        )
        result = [
            Document(
                page_content=doc["text"],
                metadata={
                    "source": file_path,
                    "file_name": os.path.basename(file_path),
                    "page_num": doc["metadata"].get("page") or "",
                },
            )
            for doc in content
        ]
        return result

    def can_handle(self, doc_type: DocumentType, content: DocumentContent) -> bool:
        return doc_type == DocumentType.PDF and content.has_text


# class DoclingLoader(BaseLoader):
#     def load(self, file_path: str) -> Dict[str, Any]:
#         converter = DocumentConverter()
#         result = converter.convert(file_path)
#         return {"content": result.document.export_to_markdown(), "format": "markdown"}

#     def can_handle(self, doc_type: DocumentType, content: DocumentContent) -> bool:
#         return doc_type == DocumentType.PDF


class UnstructuredLoader(BaseLoader):
    def load(self, file_path: str) -> Dict[str, Any]:
        elements = partition(
            filename=file_path,
            strategy="hi_res",
            include_page_breaks=True,
            skip_infer_table_types=[],
            languages=["eng", "lat"],
        )
        # content = "\n\n".join([str(el) for el in elements])
        buffer = ""
        page_numbers = []
        text = []
        current_page_num = None

        last_page_num = None
        for el in elements:
            if el.metadata.page_number is not None:
                current_page_num = el.metadata.page_number
                last_page_num = current_page_num
            else:
                current_page_num = last_page_num

            if el.category == "Table":
                dfs = pd.read_html(StringIO(str(el.metadata.text_as_html)))
                if isinstance(dfs, list):
                    df = pd.concat(dfs, ignore_index=True)
                else:
                    df = dfs.reset_index(drop=True)
                table_text = tabulate(
                    df, headers="keys", tablefmt="github", showindex=False
                )
                buffer += "\n\n" + table_text + "\n\n"
            elif el.category == "Header":
                buffer += f"\n\n## {el.text.strip()}\n\n"
            elif el.category == "Title":
                buffer += f"\n\n# {el.text.strip()}\n\n"
            elif el.category == "ListItem":
                list_text = (
                    el.text.replace("● ", "- ")
                    .replace("• ", "- ")
                    .replace("○ ", "- ")
                    .strip()
                )
                buffer += f"{list_text}\n"
            elif el.category == "PageBreak":
                text.append(buffer.strip())
                page_numbers.append(current_page_num)
                buffer = ""
            else:
                buffer += el.text + "\n"

        if buffer.strip():
            text.append(buffer.strip())
            page_numbers.append(current_page_num)

        result = [
            Document(
                page_content=content,
                metadata={
                    "source": file_path,
                    "file_name": os.path.basename(file_path),
                    "page_num": page_num,
                },
            )
            for page_num, content in zip(page_numbers, text)
        ]

        return result

    def can_handle(self, doc_type: DocumentType, content: DocumentContent) -> bool:
        return doc_type == DocumentType.PDF and content.has_text


class MarkItDownLoader(BaseLoader):
    def load(self, file_path: str) -> Dict[str, Any]:
        md = markitdown.MarkItDown()
        content = md.convert(file_path,keep_data_uris=True)
        result = Document(
            page_content=content.text_content,
            metadata={
                "source": file_path,
                "file_name": os.path.basename(file_path),
                # "page_num": page_num + 1,
            },
        )

        return result

    def can_handle(self, doc_type: DocumentType, content: DocumentContent) -> bool:
        return (
            doc_type in [DocumentType.DOCX, DocumentType.XLSX, DocumentType.XLS]
            and content.has_text
        )


class WhisperLoader(BaseLoader):
    def __init__(self):
        self.model = whisper.load_model("base")

    def load(self, file_path: str) -> Dict[str, Any]:
        content = self.model.transcribe(file_path)
        result = Document(
            page_content=content["text"],
            metadata={
                "source": file_path,
                "file_name": os.path.basename(file_path),
                "segment": content.get("segment", ""),
            },
        )

        return result

    def can_handle(self, doc_type: DocumentType, content: DocumentContent) -> bool:
        return doc_type in [
            DocumentType.MP3,
            DocumentType.WAV,
            DocumentType.M4A,
            DocumentType.MP4,
            DocumentType.AVI,
            DocumentType.MKV,
        ]


class PandasLoader(BaseLoader):
    def load(self, file_path: str) -> Dict[str, Any]:
        if file_path.endswith(".csv"):
            df = pd.read_csv(file_path)
        elif file_path.endswith((".xlsx", ".xls")):
            df = pd.read_excel(file_path)
        else:
            raise ValueError(f"Unsupported file type for pandas loader: {file_path}")

        result = Document(
            page_content=df.to_markdown(index=False),
            metadata={
                "source": file_path,
                "file_name": os.path.basename(file_path),
            },
        )
        return result

    def can_handle(self, doc_type: DocumentType, content: DocumentContent) -> bool:
        return (
            doc_type in [DocumentType.CSV, DocumentType.XLSX, DocumentType.XLS]
            and content.has_text
        )


class OpenPyXLLoader(BaseLoader):
    def load(self, file_path: str) -> Dict[str, Any]:
        wb = openpyxl.load_workbook(file_path)
        content = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            content.append(f"## {sheet_name}")
            for row in sheet.iter_rows(values_only=True):
                content.append(
                    " | ".join([str(cell) if cell is not None else "" for cell in row])
                )
        result = Document(
            page_content="\n".join(content),
            metadata={
                "source": file_path,
                "file_name": os.path.basename(file_path),
            },
        )
        return result

    def can_handle(self, doc_type: DocumentType, content: DocumentContent) -> bool:
        return doc_type in [DocumentType.XLSX, DocumentType.XLS] and content.has_text


# class PythonDocxLoader(BaseLoader):
#     def load(self, file_path: str) -> Dict[str, Any]:
#         doc = DocxDocument(file_path)
#         content = []
#         for paragraph in doc.paragraphs:
#             content.append(paragraph.text)
#         result = Document(
#             page_content="\n".join(content),
#             metadata={
#                 "source": file_path,
#                 "file_name": os.path.basename(file_path),
#             },
#         )
#         return result

#     def can_handle(self, doc_type: DocumentType, content: DocumentContent) -> bool:
#         return doc_type == DocumentType.DOCX and content.has_text


class PPTXLoader(BaseLoader):
    def load(self, file_path: str) -> Dict[str, Any]:
        prs = Presentation(file_path)
        content = []
        for i, slide in enumerate(prs.slides):
            content.append(f"## Slide {i+1}")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    content.append(shape.text)
        result = Document(
            page_content="\n".join(content),
            metadata={
                "source": file_path,
                "file_name": os.path.basename(file_path),
            },
        )
        return result

    def can_handle(self, doc_type: DocumentType, content: DocumentContent) -> bool:
        return doc_type == DocumentType.PPTX and content.has_text


class SimpleTextLoader(BaseLoader):
    def load(self, file_path: str) -> Dict[str, Any]:
        with open(file_path, "r", encoding="utf-8") as f:
            content = f.read()
        result = Document(
            page_content=content,
            metadata={
                "source": file_path,
                "file_name": os.path.basename(file_path),
            },
        )
        return result

    def can_handle(self, doc_type: DocumentType, content: DocumentContent) -> bool:
        return (
            doc_type in [DocumentType.TXT, DocumentType.JSON, DocumentType.XML]
            and content.has_text
        )
