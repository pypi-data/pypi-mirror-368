"""
Generic File System Service

This service provides a unified interface for file operations across local and S3 storage systems.
It handles various file types and returns appropriate Python objects for each format.

Features:
- Unified read/write interface for local and S3 storage
- Automatic file type detection and handling
- Support for images, documents, audio, text, and structured data
- Uses polars for CSV/dataframe operations
- Extensible file type handlers
"""

import os
import io
import tempfile
from pathlib import Path
from typing import Union, Any, Dict, BinaryIO, Optional, List, Literal, Callable, TypeVar, cast
from abc import ABC, abstractmethod
import typing
import mimetypes
from datetime import datetime, timezone

# Core dependencies - Polars disabled due to ARM64 compatibility issues
HAS_POLARS = False
pl = None

try:
    from PIL import Image
    HAS_PIL = True
except ImportError:
    HAS_PIL = False
    Image = None

# Optional dependencies - gracefully handle imports
# PDF handling is now imported from parsing.pdf_handler
try:
    from percolate.utils.parsing.pdf_handler import get_pdf_handler, HAS_PYPDF as HAS_PDF
except ImportError:
    HAS_PDF = False
    get_pdf_handler = None

try:
    import librosa
    import soundfile as sf
    HAS_AUDIO = True
except ImportError:
    HAS_AUDIO = False

try:
    import openpyxl
    import xlrd
    HAS_EXCEL = True
except ImportError:
    HAS_EXCEL = False

try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    import markdown
    HAS_MARKDOWN = True
except ImportError:
    HAS_MARKDOWN = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

from percolate.services.S3Service import S3Service
from percolate.utils import logger

# Import ResourceChunker factory function instead of direct instantiation
try:
    from percolate.utils.parsing.ResourceChunker import create_resource_chunker, get_resource_chunker
    HAS_RESOURCE_CHUNKER = True
except ImportError:
    HAS_RESOURCE_CHUNKER = False
    create_resource_chunker = None
    get_resource_chunker = None


class FileSystemProvider(ABC):
    """Abstract base class for file system providers"""
    
    @abstractmethod
    def exists(self, path: str) -> bool:
        pass
    
    @abstractmethod
    def read_bytes(self, path: str) -> bytes:
        pass
    
    @abstractmethod
    def write_bytes(self, path: str, data: bytes) -> None:
        pass
    
    @abstractmethod
    def read_text(self, path: str, encoding: str = 'utf-8') -> str:
        pass
    
    @abstractmethod
    def write_text(self, path: str, text: str, encoding: str = 'utf-8') -> None:
        pass
        
    @abstractmethod
    def open(self, path: str, mode: str = 'rb', **kwargs) -> Union[BinaryIO, Any]:
        """
        Open a file-like object.
        
        Args:
            path: The file path or URI
            mode: File mode ('r', 'rb', 'w', 'wb')
            **kwargs: Additional provider-specific arguments
            
        Returns:
            A file-like object
        """
        pass


class LocalFileSystemProvider(FileSystemProvider):
    """Local file system provider"""
    
    def _normalize_path(self, path: str) -> str:
        """Normalize path by removing file:// prefix if present"""
        if path.startswith('file://'):
            return path[7:]  # Remove 'file://' prefix
        return path
    
    def exists(self, path: str) -> bool:
        return Path(self._normalize_path(path)).exists()
    
    def read_bytes(self, path: str) -> bytes:
        normalized_path = self._normalize_path(path)
        with open(normalized_path, 'rb') as f:
            return f.read()
    
    def write_bytes(self, path: str, data: bytes) -> None:
        normalized_path = self._normalize_path(path)
        Path(normalized_path).parent.mkdir(parents=True, exist_ok=True)
        with open(normalized_path, 'wb') as f:
            f.write(data)
    
    def read_text(self, path: str, encoding: str = 'utf-8') -> str:
        normalized_path = self._normalize_path(path)
        with open(normalized_path, 'r', encoding=encoding) as f:
            return f.read()
    
    def write_text(self, path: str, text: str, encoding: str = 'utf-8') -> None:
        normalized_path = self._normalize_path(path)
        Path(normalized_path).parent.mkdir(parents=True, exist_ok=True)
        with open(normalized_path, 'w', encoding=encoding) as f:
            f.write(text)
            
    def open(self, path: str, mode: str = 'rb', **kwargs) -> BinaryIO:
        """
        Open a file-like object for local files.
        
        Args:
            path: The file path
            mode: File mode ('r', 'rb', 'w', 'wb')
            **kwargs: Additional arguments (ignored for local files)
            
        Returns:
            A file-like object
        """
        normalized_path = self._normalize_path(path)
        Path(normalized_path).parent.mkdir(parents=True, exist_ok=True)
        return open(normalized_path, mode)


class S3FileSystemProvider(FileSystemProvider):
    """S3 file system provider using existing S3Service"""
    
    def __init__(self, s3_service: Optional[S3Service] = None):
        self.s3_service = s3_service or S3Service()
    
    def exists(self, path: str) -> bool:
        try:
            parsed = self.s3_service.parse_s3_uri(path)
            bucket_name = parsed["bucket"]
            object_key = parsed["key"]
            self.s3_service.s3_client.head_object(Bucket=bucket_name, Key=object_key)
            return True
        except:
            return False
    
    def read_bytes(self, path: str) -> bytes:
        result = self.s3_service.download_file_from_uri(path)
        return result["content"]
    
    def write_bytes(self, path: str, data: bytes) -> None:
        self.s3_service.upload_filebytes_to_uri(path, data)
    
    def read_text(self, path: str, encoding: str = 'utf-8') -> str:
        data = self.read_bytes(path)
        return data.decode(encoding)
    
    def write_text(self, path: str, text: str, encoding: str = 'utf-8') -> None:
        data = text.encode(encoding)
        self.write_bytes(path, data)
        
    def open(self, path: str, mode: str = 'rb', version_id: str = None):
        """
        Open a file-like object for S3 files.
        
        Args:
            path: The S3 URI
            mode: File mode ('r', 'rb', 'w', 'wb')
            version_id: Optional version ID for versioned objects
            
        Returns:
            A file-like object
        """
        return self.s3_service.open(path, mode=mode, version_id=version_id)


class FileTypeHandler(ABC):
    """Abstract base class for file type handlers"""
    
    @abstractmethod
    def can_handle(self, file_path: str) -> bool:
        pass
    
    @abstractmethod
    def read(self, provider: FileSystemProvider, file_path: str, **kwargs) -> Any:
        pass
    
    @abstractmethod
    def write(self, provider: FileSystemProvider, file_path: str, data: Any, **kwargs) -> None:
        pass


class ImageHandler(FileTypeHandler):
    """Handler for image files (PNG, JPG, JPEG, TIFF, etc.)"""
    
    SUPPORTED_FORMATS = {'.png', '.jpg', '.jpeg', '.tiff', '.tif', '.bmp', '.gif', '.webp'}
    
    def can_handle(self, file_path: str) -> bool:
        return Path(file_path).suffix.lower() in self.SUPPORTED_FORMATS and HAS_PIL
    
    def read(self, provider: FileSystemProvider, file_path: str, **kwargs):
        if not HAS_PIL:
            raise ImportError("PIL is required for image handling but not available")
        data = provider.read_bytes(file_path)
        return Image.open(io.BytesIO(data))
    
    def write(self, provider: FileSystemProvider, file_path: str, data, **kwargs) -> None:
        if not HAS_PIL:
            raise ImportError("PIL is required for image handling but not available")
        buffer = io.BytesIO()
        format_type = kwargs.get('format', Path(file_path).suffix[1:].upper())
        if format_type.upper() == 'JPG':
            format_type = 'JPEG'
        data.save(buffer, format=format_type, **kwargs)
        provider.write_bytes(file_path, buffer.getvalue())


class TextHandler(FileTypeHandler):
    """Handler for text files (TXT, MD, HTML, etc.)"""
    
    SUPPORTED_FORMATS = {'.txt', '.md', '.html', '.htm', '.css', '.js', '.py', '.json', '.xml', '.yaml', '.yml'}
    
    def can_handle(self, file_path: str) -> bool:
        return Path(file_path).suffix.lower() in self.SUPPORTED_FORMATS
    
    def read(self, provider: FileSystemProvider, file_path: str, **kwargs) -> str:
        encoding = kwargs.get('encoding', 'utf-8')
        return provider.read_text(file_path, encoding)
    
    def write(self, provider: FileSystemProvider, file_path: str, data: str, **kwargs) -> None:
        encoding = kwargs.get('encoding', 'utf-8')
        provider.write_text(file_path, data, encoding)


class CSVHandler(FileTypeHandler):
    """Handler for CSV files using Polars"""
    
    def can_handle(self, file_path: str) -> bool:
        return Path(file_path).suffix.lower() == '.csv' and HAS_POLARS
    
    def read(self, provider: FileSystemProvider, file_path: str, **kwargs):
        if not HAS_POLARS:
            raise ImportError("polars is required for CSV handling but not available")
        data = provider.read_bytes(file_path)
        return pl.read_csv(io.BytesIO(data), **kwargs)
    
    def write(self, provider: FileSystemProvider, file_path: str, data, **kwargs) -> None:
        if not HAS_POLARS:
            raise ImportError("polars is required for CSV handling but not available")
        buffer = io.BytesIO()
        data.write_csv(buffer, **kwargs)
        provider.write_bytes(file_path, buffer.getvalue())


class ParquetHandler(FileTypeHandler):
    """Handler for Parquet files using Polars"""
    
    def can_handle(self, file_path: str) -> bool:
        return Path(file_path).suffix.lower() == '.parquet' and HAS_POLARS
    
    def read(self, provider: FileSystemProvider, file_path: str, **kwargs):
        if not HAS_POLARS:
            raise ImportError("polars is required for Parquet handling but not available")
        data = provider.read_bytes(file_path)
        return pl.read_parquet(io.BytesIO(data), **kwargs)
    
    def write(self, provider: FileSystemProvider, file_path: str, data, **kwargs) -> None:
        if not HAS_POLARS:
            raise ImportError("polars is required for Parquet handling but not available")
        buffer = io.BytesIO()
        data.write_parquet(buffer, **kwargs)
        provider.write_bytes(file_path, buffer.getvalue())


class PDFHandler(FileTypeHandler):
    """Handler for PDF files using improved implementation from pdf_handler module"""
    
    def can_handle(self, file_path: str) -> bool:
        """Check if this handler can process the file."""
        return HAS_PDF and get_pdf_handler() and get_pdf_handler().can_handle(file_path)
    
    def read(self, provider: FileSystemProvider, file_path: str, **kwargs) -> Dict[str, Any]:
        """
        Read PDF and return a dictionary with text content and metadata.
        Enhanced version that delegates to the pdf_handler module.
        """
        if not HAS_PDF:
            raise ImportError("PDF support requires pypdf")
        
        # Use the file stream from the provider
        pdf_stream = provider.open(file_path, 'rb')
        
        # Delegate to the improved PDFHandler implementation
        return get_pdf_handler().read(pdf_stream, **kwargs)
    
    def write(self, provider: FileSystemProvider, file_path: str, data: bytes, **kwargs) -> None:
        """Write PDF bytes to file"""
        provider.write_bytes(file_path, data)


class AudioHandler(FileTypeHandler):
    """Handler for audio files (WAV, MP3, etc.)"""
    
    SUPPORTED_FORMATS = {'.wav', '.mp3', '.flac', '.ogg', '.m4a'}
    
    def can_handle(self, file_path: str) -> bool:
        return Path(file_path).suffix.lower() in self.SUPPORTED_FORMATS and HAS_AUDIO
    
    def read(self, provider: FileSystemProvider, file_path: str, **kwargs) -> Dict[str, Any]:
        """Read audio file and return audio data with metadata"""
        if not HAS_AUDIO:
            raise ImportError("Audio support requires librosa and soundfile: pip install librosa soundfile")
        
        data = provider.read_bytes(file_path)
        
        # Save to temporary file for librosa (it needs file path)
        import tempfile
        with tempfile.NamedTemporaryFile(suffix=Path(file_path).suffix, delete=False) as tmp_file:
            tmp_file.write(data)
            tmp_path = tmp_file.name
        
        try:
            # Load audio data
            sr = kwargs.get('sr', None)  # Sample rate
            audio_data, sample_rate = librosa.load(tmp_path, sr=sr)
            
            # Get metadata
            info = sf.info(tmp_path)
            
            return {
                'audio_data': audio_data,
                'sample_rate': sample_rate,
                'duration': info.duration,
                'channels': info.channels,
                'format': info.format,
                'subtype': info.subtype
            }
        finally:
            os.unlink(tmp_path)  # Clean up temp file
    
    def write(self, provider: FileSystemProvider, file_path: str, data: Dict[str, Any], **kwargs) -> None:
        """Write audio data to file"""
        if not HAS_AUDIO:
            raise ImportError("Audio support requires soundfile: pip install soundfile")
        
        audio_data = data['audio_data']
        sample_rate = data['sample_rate']
        
        # Write to temporary file first
        import tempfile
        with tempfile.NamedTemporaryFile(suffix=Path(file_path).suffix, delete=False) as tmp_file:
            tmp_path = tmp_file.name
        
        try:
            sf.write(tmp_path, audio_data, sample_rate, **kwargs)
            with open(tmp_path, 'rb') as f:
                provider.write_bytes(file_path, f.read())
        finally:
            os.unlink(tmp_path)


class ExcelHandler(FileTypeHandler):
    """Handler for Excel files (XLS, XLSX) using Polars"""
    
    SUPPORTED_FORMATS = {'.xls', '.xlsx'}
    
    def can_handle(self, file_path: str) -> bool:
        return Path(file_path).suffix.lower() in self.SUPPORTED_FORMATS and HAS_EXCEL and HAS_POLARS
    
    def read(self, provider: FileSystemProvider, file_path: str, **kwargs):
        """Read Excel file and return dictionary of sheet name -> DataFrame"""
        if not HAS_EXCEL:
            raise ImportError("Excel support requires openpyxl: pip install openpyxl")
        if not HAS_POLARS:
            raise ImportError("polars is required for Excel handling but not available")
        
        data = provider.read_bytes(file_path)
        
        # Save to temporary file
        import tempfile
        with tempfile.NamedTemporaryFile(suffix=Path(file_path).suffix, delete=False) as tmp_file:
            tmp_file.write(data)
            tmp_path = tmp_file.name
        
        try:
            import pandas as pd
            
            # Filter out kwargs that aren't for pandas.read_excel
            pandas_kwargs = {k: v for k, v in kwargs.items() 
                           if k not in ['mode']}  # Remove 'mode' and other non-pandas args
            
            # Try multiple engines in order of preference for performance and reliability
            engines_to_try = []
            
            # Check if calamine is available (fastest, but newer)
            try:
                import python_calamine
                engines_to_try.append('calamine')
            except ImportError:
                pass
            
            # Always have openpyxl as fallback
            engines_to_try.append('openpyxl')
            
            excel_data = None
            last_error = None
            
            for engine in engines_to_try:
                try:
                    logger.info(f"Attempting to read Excel file with {engine} engine")
                    excel_data = pd.read_excel(tmp_path, sheet_name=None, engine=engine, **pandas_kwargs)
                    logger.info(f"Successfully read Excel file with {engine} engine")
                    break
                except Exception as e:
                    logger.warning(f"Failed to read Excel with {engine} engine: {e}")
                    last_error = e
                    continue
            
            if excel_data is None:
                raise last_error or Exception("Failed to read Excel file with any available engine")
            
            # Convert to Polars DataFrames with robust error handling
            result = {}
            for sheet_name, df in excel_data.items():
                try:
                    # First, try direct conversion
                    result[sheet_name] = pl.from_pandas(df)
                except Exception as e:
                    logger.warning(f"Direct Polars conversion failed for sheet '{sheet_name}': {e}")
                    try:
                        # Fallback 1: Handle NaN values by filling with empty strings
                        df_filled = df.fillna('')
                        result[sheet_name] = pl.from_pandas(df_filled)
                    except Exception as e2:
                        logger.warning(f"NaN-filled conversion failed for sheet '{sheet_name}': {e2}")
                        try:
                            # Fallback 2: Convert all to string to avoid PyArrow type issues
                            df_str = df.astype(str)
                            result[sheet_name] = pl.from_pandas(df_str)
                        except Exception as e3:
                            logger.error(f"All conversion methods failed for sheet '{sheet_name}': {e3}")
                            # Return the original pandas DataFrame if Polars conversion fails completely
                            result[sheet_name] = df
            
            return result
        finally:
            os.unlink(tmp_path)
    
    def write(self, provider: FileSystemProvider, file_path: str, data, **kwargs) -> None:
        """Write dictionary of DataFrames to Excel sheets"""
        if not HAS_POLARS:
            raise ImportError("polars is required for Excel handling but not available")
        import tempfile
        
        with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp_file:
            tmp_path = tmp_file.name
        
        try:
            # Convert to pandas and write
            import pandas as pd
            with pd.ExcelWriter(tmp_path, engine='openpyxl') as writer:
                for sheet_name, df in data.items():
                    pandas_df = df.to_pandas()
                    pandas_df.to_excel(writer, sheet_name=sheet_name, index=False, **kwargs)
            
            with open(tmp_path, 'rb') as f:
                provider.write_bytes(file_path, f.read())
        finally:
            os.unlink(tmp_path)


class DocxHandler(FileTypeHandler):
    """Handler for DOCX files"""
    
    def can_handle(self, file_path: str) -> bool:
        return Path(file_path).suffix.lower() == '.docx' and HAS_DOCX
    
    def read(self, provider: FileSystemProvider, file_path: str, **kwargs) -> Dict[str, Any]:
        """Read DOCX file and extract text content"""
        if not HAS_DOCX:
            raise ImportError("DOCX support requires python-docx: pip install python-docx")
        
        data = provider.read_bytes(file_path)
        
        # Save to temporary file
        import tempfile
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp_file:
            tmp_file.write(data)
            tmp_path = tmp_file.name
        
        try:
            doc = Document(tmp_path)
            
            # Extract paragraphs
            paragraphs = []
            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append(para.text)
            
            # Extract tables
            tables = []
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                tables.append(table_data)
            
            return {
                'paragraphs': paragraphs,
                'tables': tables,
                'full_text': '\n'.join(paragraphs)
            }
        finally:
            os.unlink(tmp_path)
    
    def write(self, provider: FileSystemProvider, file_path: str, data: str, **kwargs) -> None:
        """Create a simple DOCX file with the provided text"""
        if not HAS_DOCX:
            raise ImportError("DOCX support requires python-docx: pip install python-docx")
        
        import tempfile
        
        doc = Document()
        doc.add_paragraph(data)
        
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp_file:
            tmp_path = tmp_file.name
        
        try:
            doc.save(tmp_path)
            with open(tmp_path, 'rb') as f:
                provider.write_bytes(file_path, f.read())
        finally:
            os.unlink(tmp_path)


class PPTXHandler(FileTypeHandler):
    """Handler for PowerPoint presentations (.pptx, .ppt)"""
    
    def can_handle(self, file_path: str) -> bool:
        """Check if this handler can process the file"""
        ext = Path(file_path).suffix.lower()
        return ext in ['.pptx', '.ppt']
    
    def read(self, provider: FileSystemProvider, file_path: str, **kwargs) -> Dict[str, Any]:
        """
        Read PPTX file and return structured data with text and optionally images.
        
        Args:
            provider: File system provider
            file_path: Path to the PPTX file
            mode: 'simple' for text only, 'enriched' for text + LLM image analysis
            max_images: Maximum number of images to analyze (default: 50)
            
        Returns:
            Dict with slides, text_content, and optionally image_analysis
        """
        if not HAS_PPTX:
            raise ImportError("PPTX support requires python-pptx: pip install python-pptx")
        
        # Get enhanced PPTX provider from our parsing utilities
        try:
            from percolate.utils.parsing.providers import PPTXContentProvider
            pptx_provider = PPTXContentProvider()
            
            # Use temporary file for PPTX processing
            import tempfile
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp_file:
                tmp_path = tmp_file.name
                tmp_file.write(provider.read_bytes(file_path))
            
            try:
                mode = kwargs.get('mode', 'simple')
                max_images = kwargs.get('max_images', 50)
                enriched = (mode == 'enriched')
                
                # Use our enhanced provider
                text_content = pptx_provider.extract_text(tmp_path, enriched=enriched, max_images=max_images)
                
                # Also get slide-by-slide breakdown
                prs = Presentation(tmp_path)
                slides = []
                
                for slide_num, slide in enumerate(prs.slides, 1):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text.strip())
                    
                    slides.append({
                        'slide_number': slide_num,
                        'text': slide_text,
                        'combined_text': '\n'.join(slide_text)
                    })
                
                result = {
                    'slides': slides,
                    'text_content': text_content,
                    'total_slides': len(slides),
                    'mode': mode,
                    'file_path': file_path
                }
                
                # Add metadata if enriched mode and images were analyzed
                if enriched and '=== SLIDE IMAGE ANALYSIS ===' in text_content:
                    result['has_image_analysis'] = True
                    result['max_images_processed'] = max_images
                else:
                    result['has_image_analysis'] = False
                
                return result
                
            finally:
                os.unlink(tmp_path)
                
        except ImportError:
            # Fallback to basic PPTX handling without our enhanced provider
            logger.warning("Enhanced PPTX provider not available, using basic extraction")
            
            import tempfile
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp_file:
                tmp_path = tmp_file.name
                tmp_file.write(provider.read_bytes(file_path))
            
            try:
                prs = Presentation(tmp_path)
                slides = []
                all_text = []
                
                for slide_num, slide in enumerate(prs.slides, 1):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text.strip())
                    
                    slides.append({
                        'slide_number': slide_num,
                        'text': slide_text,
                        'combined_text': '\n'.join(slide_text)
                    })
                    all_text.extend(slide_text)
                
                return {
                    'slides': slides,
                    'text_content': '\n\n'.join(all_text),
                    'total_slides': len(slides),
                    'mode': 'simple',
                    'has_image_analysis': False,
                    'file_path': file_path
                }
                
            finally:
                os.unlink(tmp_path)
    
    def write(self, provider: FileSystemProvider, file_path: str, data: Any, **kwargs) -> None:
        """
        Write PPTX file (limited functionality - creates basic presentation)
        
        Args:
            provider: File system provider
            file_path: Path where to write the PPTX file
            data: Either string (single slide) or list of strings (multiple slides)
        """
        if not HAS_PPTX:
            raise ImportError("PPTX support requires python-pptx: pip install python-pptx")
        
        import tempfile
        
        prs = Presentation()
        
        if isinstance(data, str):
            # Single slide
            slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
            slide.shapes.title.text = "Generated Slide"
            slide.shapes.placeholders[1].text = data
        elif isinstance(data, list):
            # Multiple slides
            for i, slide_content in enumerate(data):
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = f"Slide {i + 1}"
                slide.shapes.placeholders[1].text = str(slide_content)
        else:
            raise ValueError(f"Unsupported data type for PPTX: {type(data)}")
        
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp_file:
            tmp_path = tmp_file.name
        
        try:
            prs.save(tmp_path)
            with open(tmp_path, 'rb') as f:
                provider.write_bytes(file_path, f.read())
        finally:
            os.unlink(tmp_path)


class FileSystemService:
    """
    Unified file system service that provides a single interface for file operations
    across local and S3 storage with automatic file type detection and handling.
    """
    
    def __init__(self, s3_service: Optional[S3Service] = None):
        self.s3_service = s3_service
        self._providers = {}
        self._handlers = []
        
        # Register default handlers
        self.register_handler(ImageHandler())
        self.register_handler(TextHandler())
        self.register_handler(CSVHandler())
        self.register_handler(ParquetHandler())
        
        if HAS_PDF:
            self.register_handler(PDFHandler())
        if HAS_AUDIO:
            self.register_handler(AudioHandler())
        if HAS_EXCEL:
            self.register_handler(ExcelHandler())
        if HAS_DOCX:
            self.register_handler(DocxHandler())
        if HAS_PPTX:
            self.register_handler(PPTXHandler())
    
    def register_handler(self, handler: FileTypeHandler):
        """Register a new file type handler"""
        self._handlers.append(handler)
    
    def _get_provider(self, path: str) -> FileSystemProvider:
        """Get the appropriate provider for the given path"""
        if path.startswith('s3://'):
            if 's3' not in self._providers:
                self._providers['s3'] = S3FileSystemProvider(self.s3_service)
            return self._providers['s3']
        else:
            if 'local' not in self._providers:
                self._providers['local'] = LocalFileSystemProvider()
            return self._providers['local']
    
    def _get_handler(self, file_path: str) -> Optional[FileTypeHandler]:
        """Get the appropriate handler for the given file path"""
        for handler in self._handlers:
            if handler.can_handle(file_path):
                return handler
        return None
    
    def _infer_file_type(self, content: bytes) -> str:
        """
        Infer file type from content by examining file signatures/magic bytes.
        Returns a file extension (with dot) that can be used for handler selection.
        """
        # Check for common file signatures
        if content.startswith(b'%PDF'):
            return '.pdf'
        elif content.startswith(b'\x89PNG\r\n\x1a\n'):
            return '.png'
        elif content.startswith(b'\xff\xd8\xff'):
            return '.jpg'
        elif content.startswith(b'GIF87a') or content.startswith(b'GIF89a'):
            return '.gif'
        elif content.startswith(b'BM'):
            return '.bmp'
        elif content.startswith(b'RIFF') and content[8:12] == b'WEBP':
            return '.webp'
        elif content.startswith(b'PK\x03\x04'):
            # ZIP-based formats - need further inspection
            if b'word/' in content[:1000]:
                return '.docx'
            elif b'xl/' in content[:1000]:
                return '.xlsx'
            elif b'ppt/' in content[:1000]:
                return '.pptx'
            else:
                return '.zip'
        elif content.startswith(b'\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1'):
            # Old Microsoft Office format
            return '.xls'  # Could also be .doc, but Excel is more common for data
        elif content.startswith(b'ID3') or content.startswith(b'\xff\xfb'):
            return '.mp3'
        elif content.startswith(b'RIFF') and content[8:12] == b'WAVE':
            return '.wav'
        elif content.startswith(b'fLaC'):
            return '.flac'
        elif content.startswith(b'OggS'):
            return '.ogg'
        elif content.startswith(b'ftypM4A'):
            return '.m4a'
        elif content.startswith((b'<!DOCTYPE', b'<!doctype', b'<html', b'<HTML')):
            return '.html'
        elif content.startswith((b'<?xml', b'<?XML')):
            return '.xml'
        elif content.startswith(b'{') or content.startswith(b'['):
            # Likely JSON, but let's check if it's valid
            try:
                content_str = content.decode('utf-8', errors='ignore')[:1000]
                import json
                json.loads(content_str)
                return '.json'
            except:
                pass
        
        # Try to detect text files by checking if content is valid UTF-8
        try:
            content.decode('utf-8')
            return '.txt'  # Default to text if UTF-8 decodable
        except:
            # Not valid UTF-8, return empty to indicate binary
            return ''
    
    def exists(self, path: str) -> bool:
        """Check if a file exists"""
        provider = self._get_provider(path)
        return provider.exists(path)
    
    def open(self, path: str, mode: str = 'rb', **kwargs) -> Union[BinaryIO, Any]:
        """
        Open a file-like object for reading or writing.
        
        This method provides a file-like interface for all supported file systems,
        allowing you to use the same code for local and S3 files.
        
        Args:
            path: File path (local or s3://)
            mode: File mode ('r', 'rb', 'w', 'wb')
            **kwargs: Additional arguments passed to the provider
                - For S3: version_id (Optional version ID for reading specific versions)
                
        Returns:
            A file-like object for reading or writing
            
        Usage:
            # Reading a file
            with fs.open("s3://bucket/key", "rb") as f:
                data = f.read()
                
            # Writing a file
            with fs.open("s3://bucket/key", "wb") as f:
                f.write(b"Hello, World!")
                
            # Reading a local file
            with fs.open("/path/to/file", "r") as f:
                text = f.read()
        """
        provider = self._get_provider(path)
        return provider.open(path, mode=mode, **kwargs)
    
    def read_bytes(self, path: str) -> bytes:
        """
        Read a file as raw bytes, bypassing any handlers.
        Useful for copying files without processing them.
        """
        provider = self._get_provider(path)
        return provider.read_bytes(path)
    
    def read(self, path: str, **kwargs) -> Any:
        """
        Read a file and return the appropriate Python object based on file type.
        
        Args:
            path: File path (local or s3://)
            **kwargs: Additional arguments passed to the specific handler
            
        Returns:
            Appropriate Python object:
            - Images: PIL.Image
            - Text files: str
            - CSV/Parquet: polars.DataFrame
            - PDF: dict with text_pages, images, metadata
            - Audio: dict with audio_data, sample_rate, metadata
            - Excel: dict of sheet_name -> polars.DataFrame
            - DOCX: dict with paragraphs, tables, full_text
            - Unknown types: bytes
        """
        logger.info(f"Reading file: {path}")
        
        provider = self._get_provider(path)
        handler = self._get_handler(path)
        
        # Special handling for audio files (WAV, MP3, etc.) when using extended mode
        ext = Path(path).suffix.lower()
        mode = kwargs.get('mode', 'simple')
        if ext in ['.wav', '.mp3', '.m4a', '.flac', '.ogg'] and mode == 'extended':
            # For audio files in extended mode, we should let ResourceChunker handle it directly
            # rather than trying to read it here, since it requires transcription
            logger.info(f"Audio file detected with extended mode: {path}")
            return provider.read_bytes(path)
        
        if handler:
            try:
                result = handler.read(provider, path, **kwargs)
                
                # Handle extended mode for PDFs
                if isinstance(handler, PDFHandler):
                    if mode == 'extended':
                        # Add raw bytes for page conversion
                        if isinstance(result, dict):
                            result['raw_bytes'] = provider.read_bytes(path)
                        
                        file_name = Path(path).name
                        extended_content = self._extract_extended_pdf_content(result, file_name, path)
                        
                        # Return enhanced result with extended content
                        if isinstance(result, dict):
                            result['content'] = extended_content
                            result['mode'] = 'extended'
                        
                        return result
                    else:
                        # Simple mode - add basic content field
                        if isinstance(result, dict):
                            text_content = '\n'.join(result.get('text_pages', []))
                            result['content'] = text_content
                            result['mode'] = 'simple'
                        
                        return result
                
                return result
            except Exception as e:
                # If handler failed for audio file, let ResourceChunker handle it
                if ext in ['.wav', '.mp3', '.m4a', '.flac', '.ogg']:
                    logger.warning(f"Audio handler failed, using raw bytes for {path}: {str(e)}")
                    return provider.read_bytes(path)
                raise
        else:
            if ext in ['.wav', '.mp3', '.m4a', '.flac', '.ogg']:
                # For audio files, return raw bytes which will be handled by ResourceChunker
                logger.info(f"No specific handler for audio file {path}. Returning raw bytes.")
                return provider.read_bytes(path)
            else:
                # No handler found - use type inference
                logger.info(f"No specific handler for {path}. Using type inference.")
                
                # Read a small sample to infer type
                sample_size = 4096  # Read first 4KB for type detection
                try:
                    with provider.open(path, 'rb') as f:
                        sample = f.read(sample_size)
                except Exception as e:
                    logger.error(f"Failed to read sample from {path}: {e}")
                    return provider.read_bytes(path)
                
                # Infer the file type
                inferred_ext = self._infer_file_type(sample)
                logger.info(f"Inferred file type for {path}: {inferred_ext or 'binary'}")
                
                if inferred_ext:
                    # Try to find a handler for the inferred type
                    pseudo_path = path + inferred_ext if not path.endswith(inferred_ext) else path
                    inferred_handler = self._get_handler(pseudo_path)
                    
                    if inferred_handler:
                        try:
                            logger.info(f"Using {inferred_handler.__class__.__name__} for inferred type {inferred_ext}")
                            return inferred_handler.read(provider, path, **kwargs)
                        except Exception as e:
                            logger.warning(f"Handler failed for inferred type {inferred_ext}: {e}")
                
                # If no handler found or handler failed, fall back to text/bytes
                if inferred_ext == '.txt':
                    try:
                        text_content = provider.read_text(path, encoding='utf-8')
                        logger.info(f"Successfully read {path} as text file ({len(text_content)} characters)")
                        return text_content
                    except Exception as text_error:
                        logger.warning(f"Failed to read {path} as text: {text_error}. Returning raw bytes.")
                        return provider.read_bytes(path)
                else:
                    # Return raw bytes for binary files
                    logger.info(f"Returning raw bytes for {path}")
                    return provider.read_bytes(path)
            
    
    def write(self, path: str, data: Any, **kwargs) -> None:
        """
        Write data to a file using the appropriate handler based on file type.
        
        Args:
            path: File path (local or s3://)
            data: Data to write (type depends on file format)
            **kwargs: Additional arguments passed to the specific handler
        """
        logger.info(f"Writing file: {path}")
        
        provider = self._get_provider(path)
        handler = self._get_handler(path)
        
        # Special case: if we're trying to write parsed PDF content,
        # we need to read the original bytes instead
        if (isinstance(data, dict) and 'text_pages' in data and 
            Path(path).suffix.lower() == '.pdf'):
            logger.warning(f"Cannot write parsed PDF content to {path}. PDF writing requires original bytes.")
            raise ValueError("Cannot write parsed PDF content. PDF files must be written as bytes.")
        
        # If we have raw bytes and a handler that expects processed data,
        # just write the bytes directly to preserve file integrity
        if isinstance(data, bytes):
            provider.write_bytes(path, data)
        elif handler:
            handler.write(provider, path, data, **kwargs)
        else:
            # Fallback: treat as text
            if isinstance(data, str):
                provider.write_text(path, data)
            else:
                raise ValueError(f"No handler for file type {Path(path).suffix} and data type {type(data)}")
    
    def copy(self, source_path: str, dest_path: str, **kwargs) -> None:
        """
        Copy a file from source to destination with optimized handling for different sources.
        
        This method intelligently handles different scenarios:
        - S3 to S3 copies: Uses direct copy operation without downloading content
        - Local to S3: Uses optimized streaming upload
        - S3 to local: Uses optimized streaming download
        - Local to local: Uses file system operations when possible
        
        Args:
            source_path: Path to the source file (local path or s3://)
            dest_path: Path to the destination file (local path or s3://)
            **kwargs: Additional arguments passed to the handlers
                - content_type: Optional MIME type for uploads
                - mode: Optional mode for handling specific file types
                
        Returns:
            None
        """
        logger.info(f"Copying file: {source_path} -> {dest_path}")
        
        source_provider = self._get_provider(source_path)
        dest_provider = self._get_provider(dest_path)
        
        # Handle S3 to S3 copy (most efficient for same bucket)
        if (
            isinstance(source_provider, S3FileSystemProvider) and 
            isinstance(dest_provider, S3FileSystemProvider)
        ):
            try:
                # Parse S3 URIs
                source_parsed = source_provider.s3_service.parse_s3_uri(source_path)
                dest_parsed = dest_provider.s3_service.parse_s3_uri(dest_path)
                
                source_bucket = source_parsed["bucket"]
                source_key = source_parsed["key"]
                dest_bucket = dest_parsed["bucket"]
                dest_key = dest_parsed["key"]
                
                # If same bucket, use copy_object for efficiency
                if source_bucket == dest_bucket:
                    logger.info(f"Using direct S3 copy_object for {source_path} -> {dest_path}")
                    source_provider.s3_service.s3_client.copy_object(
                        CopySource={'Bucket': source_bucket, 'Key': source_key},
                        Bucket=dest_bucket,
                        Key=dest_key,
                        ContentType=kwargs.get('content_type')
                    )
                    return
                else:
                    # Different buckets but still S3-to-S3, use streaming without loading into memory
                    logger.info(f"Using S3 streaming copy for {source_path} -> {dest_path}")
                    with source_provider.open(source_path, 'rb') as source_file:
                        with dest_provider.open(dest_path, 'wb') as dest_file:
                            # Stream in chunks to avoid memory issues with large files
                            chunk_size = 10 * 1024 * 1024  # 10MB chunks
                            while True:
                                chunk = source_file.read(chunk_size)
                                if not chunk:
                                    break
                                dest_file.write(chunk)
                    return
            except Exception as e:
                logger.warning(f"Direct S3 copy failed, falling back to standard copy: {str(e)}")
        
        # Handle local to local copy
        if (
            isinstance(source_provider, LocalFileSystemProvider) and 
            isinstance(dest_provider, LocalFileSystemProvider)
        ):
            try:
                import shutil
                # Normalize paths by removing file:// prefix if present
                source_normalized = source_provider._normalize_path(source_path)
                dest_normalized = dest_provider._normalize_path(dest_path)
                
                # Make sure the destination directory exists
                Path(dest_normalized).parent.mkdir(parents=True, exist_ok=True)
                
                # Use shutil for efficient local copy
                logger.info(f"Using shutil.copy for local file copy: {source_normalized} -> {dest_normalized}")
                shutil.copy2(source_normalized, dest_normalized)
                return
            except Exception as e:
                logger.warning(f"Local file copy failed, falling back to standard copy: {str(e)}")
        
        # For all other cases (or if optimized methods fail), fall back to standard copy
        logger.info(f"Using standard read/write copy for {source_path} -> {dest_path}")
        
        # Get file information to determine if special handling is needed
        source_info = self.get_file_info(source_path)
        file_extension = source_info.get('extension', '').lower()
        
        # Special handling for large files and specific formats
        if file_extension in ['.pdf', '.mp4', '.mp3', '.wav', '.zip', '.tar', '.gz']:
            # Use streaming for large files to avoid memory issues
            with source_provider.open(source_path, 'rb') as source_file:
                with dest_provider.open(dest_path, 'wb') as dest_file:
                    chunk_size = 10 * 1024 * 1024  # 10MB chunks
                    while True:
                        chunk = source_file.read(chunk_size)
                        if not chunk:
                            break
                        dest_file.write(chunk)
        else:
            # Standard copy for smaller files
            data = self.read(source_path, **kwargs)
            self.write(dest_path, data, **kwargs)
    
    def get_file_info(self, path: str) -> Dict[str, Any]:
        """Get information about a file"""
        provider = self._get_provider(path)
        
        if not provider.exists(path):
            return {'exists': False}
        
        # Get basic info
        info = {
            'exists': True,
            'path': path,
            'extension': Path(path).suffix.lower(),
            'name': Path(path).name,
            'storage_type': 's3' if path.startswith('s3://') else 'local'
        }
        
        # Try to get file size
        try:
            if path.startswith('s3://'):
                # For S3, we'd need to get object metadata
                pass  # Could implement S3 head_object here
            else:
                info['size'] = Path(path).stat().st_size
        except:
            pass
        
        # Check if we have a handler for this file type
        handler = self._get_handler(path)
        info['has_handler'] = handler is not None
        info['handler_type'] = type(handler).__name__ if handler else None
        
        # Get MIME type
        mime_type, _ = mimetypes.guess_type(path)
        info['mime_type'] = mime_type
        
        return info
    
    def apply(self, uri: str, fn: typing.Callable, **kwargs) -> typing.Any:
        """
        Apply a function to a file that requires a local file path.
        
        This method handles the case where a library function needs a file path
        (not bytes or file-like object) but we want to work with our unified
        file system interface including S3 files.
        
        Args:
            uri: File URI (local file://, S3 s3://, or HTTP/HTTPS URL)
            fn: Function that takes a file path as its first argument
            **kwargs: Additional arguments to pass to the function
            
        Returns:
            Whatever the function returns
            
        Example:
            from pdf2image import convert_from_path
            images = fs.apply("s3://bucket/document.pdf", convert_from_path, dpi=200)
        """
        import tempfile
        
        # For local files, we can pass the path directly if it's a simple file:// URI
        if uri.startswith('file://'):
            local_path = uri[7:]  # Remove 'file://' prefix
            if os.path.exists(local_path):
                return fn(local_path, **kwargs)
        elif not uri.startswith(('s3://', 'http://', 'https://')):
            # Assume it's a local path
            if os.path.exists(uri):
                return fn(uri, **kwargs)
        
        # For S3 or remote files, download to temp file
        provider = self._get_provider(uri)
        file_data = provider.read_bytes(uri)
        
        # Get file extension from URI
        file_extension = Path(uri).suffix
        
        # Create temporary file with same extension
        with tempfile.NamedTemporaryFile(
            suffix=file_extension, 
            prefix="fs_apply_", 
            delete=False
        ) as temp_file:
            temp_file.write(file_data)
            temp_file.flush()
            temp_path = temp_file.name
        
        try:
            # Apply the function to the temporary file
            result = fn(temp_path, **kwargs)
            return result
        finally:
            # Clean up temporary file
            if os.path.exists(temp_path):
                os.unlink(temp_path)
    
    def _extract_extended_pdf_content(self, pdf_data: Dict[str, Any], file_name: str, uri: str = None) -> str:
        """Extract content from PDF using LLM vision analysis of page images."""
        try:
            # Delegate to our improved PDFHandler implementation
            return get_pdf_handler().extract_extended_content(pdf_data, file_name, uri)            
        except Exception as e:
            logger.error(f"Error in extended PDF processing: {str(e)}")
            logger.info("Falling back to simple PDF parsing")
            return self._extract_simple_content(pdf_data, 'pdf')
    
    def _convert_pdf_pages_to_images(self, pdf_data: Dict[str, Any], file_name: str, uri: str = None) -> List[Any]:
        """Convert PDF pages to PIL Images for LLM analysis"""
        # Delegate to the improved PDFHandler implementation
        return get_pdf_handler().convert_pdf_to_images(pdf_data, uri)
    
    def _extract_simple_content(self, file_data: Any, file_type: str) -> str:
        """
        Extract simple text content from file data.
        
        Note: This method is kept for backward compatibility. 
        New code should use ResourceChunker handlers directly.
        """
        # Create a temporary chunker with appropriate handlers
        from percolate.utils.parsing.ResourceChunker import (
            ResourceHandler, PDFResourceHandler, TextResourceHandler,
            CSVResourceHandler, ExcelResourceHandler, DocxResourceHandler,
            PPTXResourceHandler, ImageResourceHandler, AudioResourceHandler
        )
        
        # Create handlers
        handlers = [
            PDFResourceHandler(),
            TextResourceHandler(),
            CSVResourceHandler(),
            ExcelResourceHandler(),
            DocxResourceHandler(),
            PPTXResourceHandler(),
            ImageResourceHandler(),
            AudioResourceHandler()
        ]
        
        # Find appropriate handler
        handler = None
        for h in handlers:
            if h.can_handle(file_type):
                handler = h
                break
        
        # Use handler if found, otherwise use default text handling
        if handler:
            return handler.extract_content(file_data, file_type, mode="simple")
        else:
            # Default handling for unknown types
            if isinstance(file_data, str):
                return file_data
            elif isinstance(file_data, dict):
                # Try to convert to JSON
                import json
                return json.dumps(file_data, indent=2, default=str)
            else:
                # Convert to string
                return str(file_data)
    
    def read_chunks(self, path: str, mode: str = 'simple', target_model=None, **kwargs):
        """
        Read a file and yield chunked Resources using the ResourceChunker.
        
        This is a convenient method that combines file reading and chunking in one call.
        
        Args:
            path: File path (local or s3://)
            mode: Processing mode - 'simple' (fast, text-based) or 'extended' (LLM-enhanced, expensive)
            target_model: Custom target model class for the chunks (default: Resources)
            **kwargs: Additional arguments passed to the chunker:
                - chunk_size: Maximum size of each chunk (default: 1000)
                - chunk_overlap: Overlap between chunks (default: 200) 
                - max_chunks: Maximum number of chunks to create (default: None)
                - save_to_db: Whether to save chunks to database (default: False)
                - For audio files:
                    - max_file_size_mb: Maximum file size for processing (default: 250)
                    - chunk_duration_minutes: Audio chunk duration (default: 10)
                
        Yields:
            Resources: Individual chunked Resources ready for use
            
        Example:
            # Simple chunking (fast) - iterate over chunks
            for chunk in fs.read_chunks("document.pdf"):
                print(f"Chunk: {chunk.content[:100]}...")
            
            # Extended chunking with LLM analysis (expensive but comprehensive)  
            chunks = list(fs.read_chunks("document.pdf", mode='extended', chunk_size=500))
            
            # Audio chunking with custom settings
            for chunk in fs.read_chunks("audio.wav", chunk_duration_minutes=5):
                process_audio_chunk(chunk)
            
            # Custom model with database saving
            for chunk in fs.read_chunks("spreadsheet.csv", target_model=MyCustomModel, save_to_db=True):
                handle_chunk(chunk)
        """
        logger.info(f"Reading and chunking file: {path} (mode: {mode})")
        
        # Import here to avoid circular imports
        from percolate.models.p8.types import Resources
        
        # Use custom model or default to Resources
        model_class = target_model or Resources
        
        try:
            # Get file data first for passing to the ResourceChunker
            file_data = self.read(path)
            
            # Create a ResourceChunker specifically for this FileSystemService
            if not HAS_RESOURCE_CHUNKER:
                raise ImportError("ResourceChunker is not available due to missing dependencies")
            chunker = create_resource_chunker(fs=self)
            
            # Create chunks directly using the ResourceChunker with pre-loaded data
            chunks = chunker.chunk_resource_from_uri(
                uri=path,
                parsing_mode=mode,
                chunk_size=kwargs.get('chunk_size', 1000),
                chunk_overlap=kwargs.get('chunk_overlap', 200),
                user_id=kwargs.get('userid'),
                metadata=kwargs.get('metadata'),
                file_data=file_data  # Pass the pre-loaded file data
            )
            
            # Override any additional properties if provided and the model supports them
            if chunks:
                for chunk in chunks:
                    # Override category if provided
                    if kwargs.get('category'):
                        chunk.category = kwargs['category']
                    
                    # Override name if provided
                    if kwargs.get('name'):
                        chunk_index = chunk.metadata.get('chunk_index', 0) if chunk.metadata else 0
                        if chunk_index > 0:
                            chunk.name = f"{kwargs['name']} (chunk {chunk_index + 1})"
                        else:
                            chunk.name = kwargs['name']
                    
                    # If using a custom model class, convert the chunk
                    if model_class and model_class != type(chunk):
                        # Try to create an instance of the custom model with the chunk data
                        try:
                            chunk_dict = chunk.model_dump() if hasattr(chunk, 'model_dump') else chunk.__dict__
                            # Let Pydantic handle field validation and filtering automatically
                            custom_chunk = model_class(**chunk_dict)
                            chunks[chunks.index(chunk)] = custom_chunk
                        except Exception as e:
                            logger.warning(f"Failed to convert chunk to {model_class.__name__}: {e}")
                            # Keep the original chunk if conversion fails
            
            logger.info(f"Successfully created {len(chunks)} chunks from {path}")
            
            # Yield each chunk
            for chunk in chunks:
                yield chunk
                
        except Exception as e:
            logger.error(f"Error chunking file {path}: {str(e)}")
            raise