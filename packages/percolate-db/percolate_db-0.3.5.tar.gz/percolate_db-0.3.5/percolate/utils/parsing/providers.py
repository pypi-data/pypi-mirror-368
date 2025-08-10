"""
Content providers for extracting text from various file formats.
"""
import tempfile
import os
from pathlib import Path
from abc import ABC, abstractmethod
from urllib.parse import urlparse
import requests
import fitz  # PyMuPDF
import logging

# Use standard logging instead of percolate logger to avoid circular imports
logger = logging.getLogger("percolate.parsing.providers")

# Try importing docx libraries
try:
    from docx import Document
    has_docx = True
except ImportError:
    has_docx = False
    logger.warning("python-docx not installed, DOCX support limited")

try:
    import mammoth
    has_mammoth = True
except ImportError:
    has_mammoth = False
    logger.warning("mammoth not installed, DOCX support limited")

try:
    import html2text
    has_html2text = True
except ImportError:
    has_html2text = False
    logger.warning("html2text not installed, DOCX HTML conversion not available")


def is_url(uri: str) -> bool:
    parsed = urlparse(uri)
    return parsed.scheme in ("http", "https")


def resolve_path_or_download(uri: str) -> Path:
    if Path(uri).exists():
        return Path(uri)

    if is_url(uri):
        response = requests.get(uri)
        response.raise_for_status()
        suffix = Path(urlparse(uri).path).suffix
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
        tmp.write(response.content)
        tmp.close()
        return Path(tmp.name)

    raise FileNotFoundError(f"Cannot resolve URI: {uri}")


class BaseContentProvider(ABC):
    @abstractmethod
    def extract_text(self, uri: str, enriched: bool = False) -> str:
        """Extract text from a file.
        
        Args:
            uri: File path or URL
            enriched: If True, use advanced processing (LLM analysis, etc.)
                     If False, use basic text extraction
        """
        ...


class PDFContentProvider(BaseContentProvider):
    def extract_text(self, uri: str, enriched: bool = False, max_images: int = 50) -> str:
        path = resolve_path_or_download(uri)
        
        if not enriched:
            # Raw mode: simple text extraction
            with fitz.open(str(path)) as doc:
                return "\n".join(page.get_text() for page in doc)
        else:
            # Enriched mode: extract images and use LLM analysis
            return self._extract_with_images(path, max_images)
    
    def _extract_with_images(self, path: Path, max_images: int = 50) -> str:
        """Extract text and analyze images from PDF using LLM."""
        try:
            from percolate.services.llm.ImageInterpreter import get_image_interpreter
            
            with fitz.open(str(path)) as doc:
                text_content = "\n".join(page.get_text() for page in doc)
                
                # Extract images
                images = self._extract_pdf_images(doc, max_images)
                
                if images:
                    logger.info(f"Analyzing {len(images)} images from PDF with LLM")
                    interpreter = get_image_interpreter()
                    
                    if interpreter.is_available():
                        result = interpreter.describe_images(
                            images,
                            prompt="Describe what you see in this image. Focus on text, charts, diagrams, and key visual information that would be useful for understanding the document content.",
                            context="This image was extracted from a PDF document"
                        )
                        
                        if result.get("success"):
                            image_descriptions = result["content"]
                            text_content += "\n\n=== IMAGE ANALYSIS ===\n" + str(image_descriptions)
                        else:
                            logger.warning(f"LLM image analysis failed: {result.get('error')}")
                    else:
                        logger.warning("LLM image interpreter not available, skipping image analysis")
                        
                    # Clean up temporary files
                    for img_path in images:
                        try:
                            os.unlink(img_path)
                        except:
                            pass
                
                return text_content
        except ImportError:
            logger.warning("Image interpreter not available, falling back to text extraction")
            with fitz.open(str(path)) as doc:
                return "\n".join(page.get_text() for page in doc)
    
    def _extract_pdf_images(self, doc, max_images: int = 50) -> list:
        """Extract images from PDF, sorted by size, limited to max_images."""
        images = []
        temp_dir = tempfile.mkdtemp()
        
        try:
            image_list = []
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                image_list_page = page.get_images()
                
                for img_index, img in enumerate(image_list_page):
                    try:
                        xref = img[0]
                        pix = fitz.Pixmap(doc, xref)
                        
                        if pix.n - pix.alpha < 4:  # GRAY or RGB
                            img_size = pix.width * pix.height
                            image_list.append({
                                'pix': pix,
                                'size': img_size,
                                'page': page_num,
                                'index': img_index
                            })
                        else:
                            pix = None
                    except Exception as e:
                        logger.warning(f"Failed to extract image {img_index} from page {page_num}: {e}")
            
            # Sort by size (largest first) and take top max_images
            image_list.sort(key=lambda x: x['size'], reverse=True)
            image_list = image_list[:max_images]
            
            # Save sorted images
            for i, img_data in enumerate(image_list):
                try:
                    pix = img_data['pix']
                    img_path = os.path.join(temp_dir, f"pdf_image_{i:03d}.png")
                    pix.save(img_path)
                    images.append(img_path)
                    pix = None
                except Exception as e:
                    logger.warning(f"Failed to save image {i}: {e}")
            
        except Exception as e:
            logger.error(f"Error extracting PDF images: {e}")
        
        return images


class DefaultContentProvider(BaseContentProvider):
    def extract_text(self, uri: str, enriched: bool = False) -> str:
        path = resolve_path_or_download(uri)
        text = path.read_text()
        
        if enriched:
            # Enriched mode: TODO - use LLM to interpret/summarize content
            logger.warning("Default enriched mode not yet implemented, falling back to raw")
        
        return text


class DOCXContentProvider(BaseContentProvider):
    """Content provider for Microsoft Word documents."""
    
    def extract_text(self, uri: str, enriched: bool = False) -> str:
        """Extract text from a DOCX file."""
        path = resolve_path_or_download(uri)
        
        # If no libraries available, fall back to basic text extraction
        if not has_docx and not has_mammoth:
            logger.warning("No DOCX libraries available, falling back to simple text extraction")
            return path.read_text(errors='ignore')
        
        try:
            # First try with python-docx for simple text extraction
            if has_docx:
                doc = Document(str(path))
                paragraphs = []
                for paragraph in doc.paragraphs:
                    if paragraph.text.strip():
                        paragraphs.append(paragraph.text)
                
                # If we got text, return it
                if paragraphs:
                    return '\n\n'.join(paragraphs)
            
            # Fallback to mammoth for more complex documents
            if has_mammoth:
                logger.info("Using mammoth for DOCX extraction")
                with open(str(path), "rb") as docx_file:
                    result = mammoth.convert_to_markdown(docx_file)
                    
                    if result.messages:
                        for message in result.messages:
                            logger.warning(f"DOCX conversion warning: {message}")
                    
                    return result.value
                    
        except Exception as e:
            logger.error(f"Error extracting text from DOCX: {e}")
            # Final fallback - try mammoth HTML conversion if available
            if has_mammoth and has_html2text:
                try:
                    with open(str(path), "rb") as docx_file:
                        result = mammoth.convert_to_html(docx_file)
                        h = html2text.HTML2Text()
                        h.ignore_links = False
                        return h.handle(result.value)
                except Exception as e2:
                    logger.error(f"Failed all DOCX extraction methods: {e2}")
            
            # Ultimate fallback
            return path.read_text(errors='ignore')


# New provider classes for additional formats
class HTMLContentProvider(BaseContentProvider):
    def extract_text(self, uri: str, enriched: bool = False) -> str:
        path = resolve_path_or_download(uri)
        
        if not enriched:
            # Raw mode: strip HTML tags
            try:
                if has_html2text:
                    h = html2text.HTML2Text()
                    h.ignore_links = True
                    return h.handle(path.read_text())
                else:
                    import re
                    html_content = path.read_text()
                    return re.sub(r'<[^>]+>', '', html_content)
            except Exception as e:
                logger.error(f"HTML parsing failed: {e}")
                return path.read_text()
        else:
            # Enriched mode: TODO - use LLM to analyze structure and content
            logger.warning("HTML enriched mode not yet implemented, falling back to raw")
            return self.extract_text(uri, enriched=False)


class MarkdownContentProvider(BaseContentProvider):
    def extract_text(self, uri: str, enriched: bool = False) -> str:
        path = resolve_path_or_download(uri)
        text = path.read_text()
        
        if enriched:
            # Enriched mode: TODO - parse markdown structure, extract metadata
            logger.warning("Markdown enriched mode not yet implemented, falling back to raw")
        
        return text


class XLSXContentProvider(BaseContentProvider):
    def extract_text(self, uri: str, enriched: bool = False) -> str:
        path = resolve_path_or_download(uri)
        
        try:
            import pandas as pd
            df = pd.read_excel(str(path), sheet_name=None)  # Read all sheets
            
            if not enriched:
                # Raw mode: concatenate all cell values
                all_text = []
                for sheet_name, sheet_df in df.items():
                    all_text.append(f"Sheet: {sheet_name}\n")
                    all_text.append(sheet_df.to_string(index=False))
                return "\n\n".join(all_text)
            else:
                # Enriched mode: TODO - analyze data structure, detect patterns
                logger.warning("XLSX enriched mode not yet implemented, falling back to raw")
                return self.extract_text(uri, enriched=False)
                
        except ImportError:
            logger.warning("pandas not available, falling back to default provider")
            return DefaultContentProvider().extract_text(uri, enriched=enriched)
        except Exception as e:
            logger.error(f"XLSX parsing failed: {e}")
            return f"Error reading XLSX file: {e}"


class AudioContentProvider(BaseContentProvider):
    """Content provider for audio files (WAV, MP3, etc.)."""
    
    def extract_text(self, uri: str, enriched: bool = False) -> str:
        path = resolve_path_or_download(uri)
        suffix = path.suffix.lower()
        
        if not enriched:
            # Raw mode: basic file info
            try:
                # Handle WAV files - we can get detailed info
                if suffix == '.wav':
                    import wave
                    with wave.open(str(path), 'rb') as wav_file:
                        frames = wav_file.getnframes()
                        sample_rate = wav_file.getframerate()
                        duration = frames / sample_rate
                        return f"Audio file: {path.name}\nDuration: {duration:.2f} seconds\nSample rate: {sample_rate} Hz\nFrames: {frames}"
                
                # For other audio types, just provide basic file info
                else:
                    file_size = os.path.getsize(path)
                    return f"Audio file: {path.name}\nSize: {file_size/1024:.1f} KB\nFormat: {suffix[1:]}"
                    
            except Exception as e:
                return f"Audio file: {path.name}\nError reading file: {e}"
        else:
            # Enriched mode: Audio files should be processed through ResourceChunker's media processing pipeline
            logger.warning(f"Audio file {path.name} should be processed through ResourceChunker's media processing pipeline")
            return f"Audio file: {path.name}\nTranscription requires ResourceChunker media processing pipeline"


class PPTXContentProvider(BaseContentProvider):
    def extract_text(self, uri: str, enriched: bool = False, max_images: int = 50) -> str:
        path = resolve_path_or_download(uri)
        
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            prs = Presentation(str(path))
            
            if not enriched:
                # Raw mode: extract all text from slides
                text_runs = []
                for slide_num, slide in enumerate(prs.slides, 1):
                    text_runs.append(f"Slide {slide_num}:")
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_runs.append(shape.text)
                return "\n\n".join(text_runs)
            else:
                # Enriched mode: analyze slide structure, extract and analyze images
                return self._extract_with_images(prs, max_images)
                
        except ImportError:
            logger.warning("python-pptx not available, falling back to default provider")
            return DefaultContentProvider().extract_text(uri, enriched=enriched)
        except Exception as e:
            logger.error(f"PPTX parsing failed: {e}")
            return f"Error reading PPTX file: {e}"
    
    def _extract_with_images(self, prs, max_images: int = 50) -> str:
        """Extract text and analyze images from PPTX using LLM."""
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            from percolate.services.llm.ImageInterpreter import get_image_interpreter
            
            # Extract text first
            text_runs = []
            for slide_num, slide in enumerate(prs.slides, 1):
                text_runs.append(f"Slide {slide_num}:")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_runs.append(shape.text)
            
            text_content = "\n\n".join(text_runs)
            
            # Extract images
            images = self._extract_pptx_images(prs, max_images)
            
            if images:
                logger.info(f"Analyzing {len(images)} images from PPTX with LLM")
                interpreter = get_image_interpreter()
                
                if interpreter.is_available():
                    result = interpreter.describe_images(
                        images,
                        prompt="Describe what you see in this image from a presentation slide. Focus on charts, diagrams, visual data, and any text that might not be captured as regular slide text.",
                        context="This image was extracted from a PowerPoint presentation slide"
                    )
                    
                    if result.get("success"):
                        image_descriptions = result["content"]
                        text_content += "\n\n=== SLIDE IMAGE ANALYSIS ===\n" + str(image_descriptions)
                    else:
                        logger.warning(f"LLM image analysis failed: {result.get('error')}")
                else:
                    logger.warning("LLM image interpreter not available, skipping image analysis")
                    
                # Clean up temporary files
                for img_path in images:
                    try:
                        os.unlink(img_path)
                    except:
                        pass
            
            return text_content
            
        except ImportError:
            logger.warning("Image interpreter not available, falling back to text extraction")
            # Fall back to raw mode
            text_runs = []
            for slide_num, slide in enumerate(prs.slides, 1):
                text_runs.append(f"Slide {slide_num}:")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            return "\n\n".join(text_runs)
    
    def _extract_pptx_images(self, prs, max_images: int = 50) -> list:
        """Extract images from PPTX, sorted by size, limited to max_images."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        
        images = []
        temp_dir = tempfile.mkdtemp()
        image_data = []
        
        try:
            # First pass: collect all images with size info
            for slide_number, slide in enumerate(prs.slides, start=1):
                for shape_idx, shape in enumerate(slide.shapes):
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            image = shape.image
                            image_bytes = image.blob
                            image_size = len(image_bytes)
                            image_ext = image.ext or 'png'
                            
                            image_data.append({
                                'bytes': image_bytes,
                                'size': image_size,
                                'ext': image_ext,
                                'slide': slide_number,
                                'shape': shape_idx
                            })
                        except Exception as e:
                            logger.warning(f"Failed to extract image from slide {slide_number}, shape {shape_idx}: {e}")
            
            # Sort by size (largest first) and take top max_images
            image_data.sort(key=lambda x: x['size'], reverse=True)
            image_data = image_data[:max_images]
            
            # Second pass: save the selected images
            for i, img_info in enumerate(image_data):
                try:
                    filename = f"pptx_image_{i:03d}_slide{img_info['slide']}.{img_info['ext']}"
                    img_path = os.path.join(temp_dir, filename)
                    
                    with open(img_path, 'wb') as f:
                        f.write(img_info['bytes'])
                    
                    images.append(img_path)
                    
                except Exception as e:
                    logger.warning(f"Failed to save image {i}: {e}")
                    
        except Exception as e:
            logger.error(f"Error extracting PPTX images: {e}")
        
        return images


content_providers = {
    ".pdf": PDFContentProvider(),
    ".docx": DOCXContentProvider(),
    ".doc": DOCXContentProvider(),  # Will handle old doc format too
    ".txt": DefaultContentProvider(),
    ".html": HTMLContentProvider(),
    ".htm": HTMLContentProvider(),
    ".md": MarkdownContentProvider(),
    ".markdown": MarkdownContentProvider(),
    ".xlsx": XLSXContentProvider(),
    ".xls": XLSXContentProvider(),
    ".wav": AudioContentProvider(),
    ".mp3": AudioContentProvider(),
    ".m4a": AudioContentProvider(),
    ".flac": AudioContentProvider(),
    ".ogg": AudioContentProvider(),
    ".pptx": PPTXContentProvider(),
    ".ppt": PPTXContentProvider(),
}

default_provider = DefaultContentProvider()


def get_content_provider_for_uri(uri: str) -> BaseContentProvider:
    """Get the appropriate content provider for a given URI."""
    suffix = Path(urlparse(uri).path).suffix.lower()
    return content_providers.get(suffix, default_provider)