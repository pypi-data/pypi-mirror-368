#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
docsray/scripts/file_converter.py
Convert various file formats to PDF for processing
"""

import os
import sys
from pathlib import Path
from typing import Optional, Tuple

import pypandoc

from PIL import Image

import pdfkit

import markdown

from llama_index.readers.file import HWPReader

from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4
from PIL import Image
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch    

from bs4 import BeautifulSoup


from typing import List, Tuple
from pathlib import Path
from docsray.config import CACHE_DIR
import pandas as pd

from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.lib.utils import ImageReader
import os

# Import multimedia processing utilities
try:
    from docsray.multimedia_processor import process_audio_file, process_video_file
    MULTIMEDIA_SUPPORT = True
except ImportError:
    MULTIMEDIA_SUPPORT = False
    print("Warning: Multimedia support not available. Install whisper and ffmpeg for audio/video processing.", file=sys.stderr)

def _save_text_images_to_pdf_korean(text: str, image_paths: List[str], output_file: Path) -> bool:
    """Save text and images to PDF with Korean font support"""
    try:          
        # Register Korean fonts
        # Try to find Korean fonts on the system
        korean_fonts = [
            # Windows
            "C:/Windows/Fonts/malgun.ttf",  # Malgun Gothic
            "C:/Windows/Fonts/malgunbd.ttf",  # Malgun Gothic Bold
            "C:/Windows/Fonts/gulim.ttc",   # Gulim
            "C:/Windows/Fonts/batang.ttc",  # Batang
            "C:/Windows/Fonts/NanumGothic.ttf",
            # macOS
            "/System/Library/Fonts/AppleSDGothicNeo.ttc",
            "/System/Library/Fonts/AppleGothic.ttf",
            "/Library/Fonts/NanumGothic.ttf",
            "/System/Library/Fonts/Helvetica.ttc",  # Fallback
            # Linux
            "/usr/share/fonts/truetype/nanum/NanumGothic.ttf",
            "/usr/share/fonts/truetype/fonts-nanum/NanumGothic.ttf",
            "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",  # Fallback
        ]
        
        font_registered = False
        font_name = "KoreanFont"
        
        for font_path in korean_fonts:
            if os.path.exists(font_path):
                try:
                    pdfmetrics.registerFont(TTFont(font_name, font_path))
                    font_registered = True
                    print(f"Registered Korean font: {font_path}", file=sys.stderr)
                    break
                except Exception as e:
                    print(f"Failed to register font {font_path}: {e}", file=sys.stderr)
                    continue
        
        if not font_registered:
            try:
                # Fallback: download and use a free Korean font
                import urllib.request
                font_url = "https://github.com/google/fonts/raw/main/ofl/nanumgothic/NanumGothic-Regular.ttf"
                font_dir = output_file.parent / "fonts"
                font_dir.mkdir(exist_ok=True)
                font_path = font_dir / "NanumGothic.ttf"
                
                if not font_path.exists():
                    print(f"Downloading Korean font...", file=sys.stderr)
                    urllib.request.urlretrieve(font_url, font_path)
                
                pdfmetrics.registerFont(TTFont(font_name, str(font_path)))
                font_registered = True
                print(f"Downloaded and registered Korean font", file=sys.stderr)
            except Exception as e:
                print(f"Failed to download Korean font: {e}", file=sys.stderr)
                # Use default font
                font_name = "Helvetica"
                font_registered = True
        
        if not font_registered:
            # Last resort: use Helvetica
            font_name = "Helvetica"
            print(f"Using fallback font: Helvetica", file=sys.stderr)
        
        # Create PDF
        c = canvas.Canvas(str(output_file), pagesize=A4)
        width, height = A4
        
        # Set font with Korean support
        c.setFont(font_name, 12)
        
        # Add text with better handling
        y = height - 50
        line_height = 18
        max_chars_per_line = 70
        
        for line in text.split('\n'):
            if y < 50:  # New page if needed
                c.showPage()
                c.setFont(font_name, 12)
                y = height - 50
            
            # Skip empty lines but add some space
            if not line.strip():
                y -= line_height // 2
                continue
            
            # Handle long lines by wrapping
            if len(line) > max_chars_per_line:
                words = line.split()
                current_line = ""
                for word in words:
                    test_line = current_line + (" " if current_line else "") + word
                    if len(test_line) <= max_chars_per_line:
                        current_line = test_line
                    else:
                        if current_line:
                            try:
                                c.drawString(50, y, current_line)
                            except Exception as e:
                                # Handle text encoding issues
                                safe_text = current_line.encode('utf-8', errors='replace').decode('utf-8')
                                c.drawString(50, y, safe_text)
                            y -= line_height
                            if y < 50:
                                c.showPage()
                                c.setFont(font_name, 12)
                                y = height - 50
                        current_line = word
                
                if current_line:
                    try:
                        c.drawString(50, y, current_line)
                    except Exception as e:
                        safe_text = current_line.encode('utf-8', errors='replace').decode('utf-8')
                        c.drawString(50, y, safe_text)
                    y -= line_height
            else:
                try:
                    c.drawString(50, y, line)
                except Exception as e:
                    # Handle text encoding issues
                    safe_text = line.encode('utf-8', errors='replace').decode('utf-8')
                    c.drawString(50, y, safe_text)
                y -= line_height
        
        # Add images with better error handling
        for i, img_path in enumerate(image_paths):
            if os.path.exists(img_path):
                try:
                    c.showPage()
                    img = ImageReader(img_path)
                    
                    # Calculate image dimensions to fit page
                    max_width = width - 100
                    max_height = height - 100
                    
                    c.drawImage(img, 50, 50, 
                              width=max_width, height=max_height, 
                              preserveAspectRatio=True, mask='auto')
                    
                    print(f"Added image {i+1}/{len(image_paths)}: {os.path.basename(img_path)}", file=sys.stderr)
                except Exception as e:
                    print(f"Failed to add image {img_path}: {e}", file=sys.stderr)
                    # Add a placeholder text for failed images
                    try:
                        c.showPage()
                        c.setFont(font_name, 12)
                        c.drawString(50, height/2, f"[Image could not be displayed: {os.path.basename(img_path)}]")
                    except:
                        pass
        
        c.save()
        return True
        
    except Exception as e:
        print(f"PDF creation error: {e}", file=sys.stderr)
        import traceback
        traceback.print_exc(file=sys.stderr)
        return False
            
class FileConverter:
    """Convert various file formats to PDF"""
    
    SUPPORTED_FORMATS = {
        # Office documents
        '.docx': 'Microsoft Word',
        '.doc': 'Microsoft Word (Legacy)',
        '.xlsx': 'Microsoft Excel',
        '.xls': 'Microsoft Excel (Legacy)',
        '.pdf': 'PDF Document',
        '.pptx': 'Microsoft PowerPoint',
        '.ppt': 'Microsoft PowerPoint (Legacy)',
        '.odt': 'OpenDocument Text',
        '.ods': 'OpenDocument Spreadsheet',
        '.odp': 'OpenDocument Presentation',
        '.hwp': 'Hancom Word Processor',
        '.hwpx': 'Hancom Word Processor (OOXML)',

        # Text formats
        '.txt': 'Plain Text',
        
        # Image formats
        '.jpg': 'JPEG Image',
        '.jpeg': 'JPEG Image',
        '.jfif': 'JPEG Image',
        '.png': 'PNG Image',
        '.gif': 'GIF Image',
        '.bmp': 'Bitmap Image',
        '.tiff': 'TIFF Image',
        '.tif': 'TIFF Image',
        '.webp': 'WebP Image',
        
        # Audio formats
        '.mp3': 'MP3 Audio',
        '.wav': 'WAV Audio',
        '.m4a': 'M4A Audio',
        '.flac': 'FLAC Audio',
        '.ogg': 'OGG Audio',
        '.wma': 'WMA Audio',
        '.aac': 'AAC Audio',
        
        # Video formats
        '.mp4': 'MP4 Video',
        '.avi': 'AVI Video',
        '.mov': 'MOV Video',
        '.wmv': 'WMV Video',
        '.flv': 'FLV Video',
        '.mkv': 'MKV Video',
        '.webm': 'WebM Video',
        '.m4v': 'M4V Video',
        '.mpg': 'MPEG Video',
        '.mpeg': 'MPEG Video',
    }
    
    def __init__(self, output_dir: Optional[Path] = None):
        """
        Initialize converter with optional output directory
        
        Args:
            output_dir: Directory to save converted PDFs (default: temp directory)
        """
        self.output_dir = output_dir or CACHE_DIR
        self.output_dir.mkdir(parents=True, exist_ok=True)
        
        # Check LibreOffice availability on initialization
        self._libreoffice_available = self._check_libreoffice()
        if self._libreoffice_available:
            print("LibreOffice detected - will use for office document conversion", file=sys.stderr)
        
    @classmethod
    def get_supported_formats(cls) -> dict:
        """Get dictionary of supported formats"""
        return cls.SUPPORTED_FORMATS.copy()
    
    @classmethod
    def is_supported(cls, file_path: str) -> bool:
        """Check if file format is supported"""
        ext = Path(file_path).suffix.lower()
        return ext in cls.SUPPORTED_FORMATS or ext == '.pdf'

    def convert_to_pdf(self, input_path: str, output_path: Optional[str] = None) -> Tuple[bool, str]:
        """
        Convert file to PDF
        
        Args:
            input_path: Path to input file
            output_path: Optional output path (default: auto-generated)
            
        Returns:
            Tuple of (success: bool, output_path_or_error: str)
        """
        input_file = Path(input_path)
        
        # Check if file exists
        if not input_file.exists():
            return False, f"File not found: {input_path}"
        
        # Check if already PDF
        if input_file.suffix.lower() == '.pdf':
            return True, str(input_file)
        
        # Check if format is supported
        file_ext = input_file.suffix.lower()
        if file_ext not in self.SUPPORTED_FORMATS:
            return False, f"Unsupported format: {file_ext}"
        
        # Generate output path if not provided
        if output_path is None:
            output_path = self.output_dir / f"{input_file.stem}_converted.pdf"
        else:
            output_path = Path(output_path)
        
        # Select conversion method based on file type
        print(f"Converting {self.SUPPORTED_FORMATS[file_ext]} file to PDF...", file=sys.stderr)
        
        # Office documents
        if file_ext in ['.docx', '.doc']:
            return self._convert_docx_to_pdf(input_file, output_path)
        elif file_ext in ['.xlsx', '.xls']:
            return self._convert_excel_to_pdf(input_file, output_path)
        elif file_ext in ['.pptx', '.ppt']:
            return self._convert_ppt_to_pdf(input_file, output_path)
        elif file_ext in ['.odt', '.ods', '.odp']:
            # OpenDocument formats - use LibreOffice if available, otherwise try pandoc
            return self._convert_opendocument_to_pdf(input_file, output_path)
        elif file_ext == '.hwp':
            return self._convert_hwp_to_pdf(input_file, output_path)
        elif file_ext == '.hwpx':
            return self._convert_hwpx_to_pdf(input_file, output_path)
        
        elif file_ext in ['.pdf']:
            # If already PDF, just return the path
            return True, str(input_file)

        # Text formats
        elif file_ext == '.txt':
            return self._convert_text_to_pdf(input_file, output_path)
        
        # Image formats
        elif file_ext in ['.jpg', '.jpeg', '.jfif', '.png', '.gif', '.bmp', '.tiff', '.tif', '.webp']:
            return self._convert_image_to_pdf(input_file, output_path)
        
        # Audio formats
        elif file_ext in ['.mp3', '.wav', '.m4a', '.flac', '.ogg', '.wma', '.aac']:
            return self._convert_audio_to_pdf(input_file, output_path)
        
        # Video formats
        elif file_ext in ['.mp4', '.avi', '.mov', '.wmv', '.flv', '.mkv', '.webm', '.m4v', '.mpg', '.mpeg']:
            return self._convert_video_to_pdf(input_file, output_path)
        
        # Fallback: try pandoc for anything else
        else:
            return self._convert_with_pandoc_simple(input_file, output_path)
    
    def _check_libreoffice(self) -> bool:
        """Check if LibreOffice is installed and available"""
        import subprocess
        
        # Try different possible commands
        commands = ['libreoffice', 'soffice']
        
        for cmd in commands:
            try:
                result = subprocess.run(
                    [cmd, '--version'],
                    capture_output=True,
                    check=True,
                    timeout=5
                )
                if result.returncode == 0:
                    return True
            except (subprocess.CalledProcessError, FileNotFoundError, subprocess.TimeoutExpired):
                continue
        
        return False
    
    def _convert_with_libreoffice(self, input_file: Path, output_file: Path) -> Tuple[bool, str]:
        """Convert document to PDF using LibreOffice headless mode"""
        import subprocess
        import shutil
        
        # Find the correct LibreOffice command
        libreoffice_cmd = None
        for cmd in ['libreoffice', 'soffice']:
            if shutil.which(cmd):
                libreoffice_cmd = cmd
                break
        
        if not libreoffice_cmd:
            return False, "LibreOffice not found"
        
        try:
            # Create output directory if it doesn't exist
            output_file.parent.mkdir(parents=True, exist_ok=True)
            
            # LibreOffice outputs to the directory with the original filename + .pdf
            # So we need to use a temp directory and move the file
            temp_output_dir = output_file.parent / f"libreoffice_temp_{os.getpid()}"
            temp_output_dir.mkdir(exist_ok=True)
            
            cmd = [
                libreoffice_cmd,
                '--headless',
                '--invisible',
                '--nodefault',
                '--nolockcheck',
                '--nologo',
                '--norestore',
                '--convert-to', 'pdf',
                '--outdir', str(temp_output_dir),
                str(input_file)
            ]
            
            print(f"Running LibreOffice conversion: {' '.join(cmd)}", file=sys.stderr)
            
            result = subprocess.run(
                cmd,
                capture_output=True,
                text=True,
                timeout=120  # 2 minute timeout
            )
            
            if result.returncode == 0:
                # Find the output file (LibreOffice creates it with the input filename + .pdf)
                expected_output = temp_output_dir / f"{input_file.stem}.pdf"
                
                if expected_output.exists():
                    # Move to desired location
                    shutil.move(str(expected_output), str(output_file))
                    
                    # Clean up temp directory
                    try:
                        shutil.rmtree(temp_output_dir)
                    except:
                        pass
                    
                    print(f"LibreOffice conversion successful: {output_file}", file=sys.stderr)
                    return True, str(output_file)
                else:
                    return False, f"LibreOffice did not create expected output file: {expected_output}"
            else:
                error_msg = f"LibreOffice conversion failed with code {result.returncode}"
                if result.stderr:
                    error_msg += f": {result.stderr}"
                return False, error_msg
                
        except subprocess.TimeoutExpired:
            return False, "LibreOffice conversion timed out"
        except Exception as e:
            return False, f"LibreOffice conversion error: {str(e)}"
        finally:
            # Clean up temp directory if it still exists
            if 'temp_output_dir' in locals() and temp_output_dir.exists():
                try:
                    shutil.rmtree(temp_output_dir)
                except:
                    pass
    
    def _convert_opendocument_to_pdf(self, input_file: Path, output_file: Path) -> Tuple[bool, str]:
        """Convert OpenDocument formats (ODT, ODS, ODP) to PDF"""
        # Try LibreOffice first if available (best for OpenDocument)
        if self._libreoffice_available:
            print(f"Using LibreOffice for {input_file.suffix} conversion...", file=sys.stderr)
            success, result = self._convert_with_libreoffice(input_file, output_file)
            if success:
                return success, result
            print(f"LibreOffice conversion failed: {result}", file=sys.stderr)
        
        # Fallback to pandoc for ODT files
        if input_file.suffix.lower() == '.odt':
            print(f"Trying pandoc for ODT conversion...", file=sys.stderr)
            return self._convert_with_pandoc_simple(input_file, output_file)
        
        # For ODS and ODP without LibreOffice, we can't convert easily
        return False, f"Cannot convert {input_file.suffix} files without LibreOffice installed"
        
    # ------------------------------------------------------------------
    # HWPX → PDF  (zip + XML)  — text + images via BeautifulSoup + ReportLab
    # ------------------------------------------------------------------
    def _convert_hwpx_to_pdf(self, input_file: Path, output_file: Path) -> Tuple[bool, str]:
        """Convert HWPX (Hangul OOXML) to PDF."""
        # Try LibreOffice first if available
        if self._libreoffice_available:
            print(f"Trying LibreOffice for HWPX conversion...", file=sys.stderr)
            success, result = self._convert_with_libreoffice(input_file, output_file)
            if success:
                return success, result
            print(f"LibreOffice conversion failed, falling back to manual parsing: {result}", file=sys.stderr)
        
        # Fallback to manual HWPX parsing
        from zipfile import ZipFile
        texts, img_paths = [], []
        tmp_dir = output_file.with_suffix("")  # temp folder alongside output
        tmp_dir.mkdir(exist_ok=True)

        try:
            with ZipFile(input_file) as zf:
                # 1) Collect text + image ids
                for name in zf.namelist():
                    if name.startswith("Contents/") and name.endswith(".xml"):
                        # Read with proper encoding
                        raw_content = zf.read(name)
                        
                        # Try different encodings for Korean text
                        for encoding in ['utf-8', 'utf-8-sig', 'cp949', 'euc-kr']:
                            try:
                                content = raw_content.decode(encoding)
                                break
                            except UnicodeDecodeError:
                                continue
                        else:
                            # If all encodings fail, use utf-8 with error handling
                            content = raw_content.decode('utf-8', errors='replace')
                        
                        soup = BeautifulSoup(content, "xml")
                        
                        # HWPX uses different tag structure
                        # Look for text in various possible tags
                        text_tags = ['t', 'text', 'hp:t', 'hp:run', 'w:t']
                        for tag_name in text_tags:
                            for text_elem in soup.find_all(tag_name):
                                text = text_elem.get_text(strip=True)
                                if text:
                                    texts.append(text)
                        
                        # Also try to find paragraphs with different naming conventions
                        para_tags = ['p', 'para', 'hp:p', 'hp:para', 'w:p']
                        for tag_name in para_tags:
                            for para in soup.find_all(tag_name):
                                para_text = para.get_text(" ", strip=True)
                                if para_text and para_text not in texts:
                                    texts.append(para_text)
                        
                        # Find images with various possible tag/attribute names
                        img_tags = ['pic', 'hp:pic', 'image', 'hp:image']
                        img_attrs = ['binFile', 'hp:binFile', 'r:embed', 'href']
                        
                        for img_tag in img_tags:
                            for pic in soup.find_all(img_tag):
                                for attr in img_attrs:
                                    img_id = pic.get(attr)
                                    if img_id:
                                        img_paths.append(img_id)

                # Debug: Print what we found
                print(f"Found {len(texts)} text elements", file=sys.stderr)
                print(f"Found {len(img_paths)} image references", file=sys.stderr)
                
                # 2) Extract BinData images
                extracted_imgs = []
                for img_id in img_paths:
                    for possible_path in [f"BinData/{img_id}", f"Contents/BinData/{img_id}", img_id]:
                        try:
                            data = zf.read(possible_path)
                            # Detect image type
                            if data[:8] == b"\x89PNG\r\n\x1a\n":
                                ext = ".png"
                            elif data[:2] == b"\xff\xd8":
                                ext = ".jpg"
                            else:
                                ext = ".bin"  # unknown format
                            
                            out_path = tmp_dir / f"{input_file.stem}_{img_id}{ext}"
                            out_path.write_bytes(data)
                            extracted_imgs.append(str(out_path))
                            break
                        except KeyError:
                            continue

                # 3) Synthesize PDF with proper encoding
                combined_text = "\n".join(texts) if texts else "No text content found"
                
                # Ensure _save_text_images_to_pdf handles Korean text properly
                success = _save_text_images_to_pdf_korean(combined_text, extracted_imgs, output_file)
                if success:
                    return True, str(output_file)
                return False, "PDF synthesis failed"

        except Exception as e:
            return False, f"HWPX conversion error: {e}"
        finally:
            # Cleanup temp directory
            try:
                import shutil
                shutil.rmtree(tmp_dir)
            except:
                pass

    def _convert_hwp_to_pdf(self, input_file: Path, output_file: Path) -> Tuple[bool, str]:
        """Convert HWP to PDF"""
        # Try LibreOffice first if available
        if self._libreoffice_available:
            print(f"Trying LibreOffice for HWP conversion...", file=sys.stderr)
            success, result = self._convert_with_libreoffice(input_file, output_file)
            if success:
                return success, result
            print(f"LibreOffice conversion failed, falling back to HWPReader: {result}", file=sys.stderr)
        
        # Fallback to HWPReader
        try:
            reader = HWPReader()
            # llama_index HWPReader returns a list of Document objects
            documents = reader.load_data(file=str(input_file))
            
            if documents and len(documents) > 0:
                # Combine text from all documents
                full_text = ""
                all_images = []
                
                for doc in documents:
                    # llama_index Document objects have a 'text' attribute
                    full_text += doc.text + "\n\n"
                    
                    # Extract images from metadata
                    # Metadata structure may vary depending on HWPReader implementation
                    if hasattr(doc, 'metadata') and doc.metadata:
                        images = doc.metadata.get("images", [])
                        all_images.extend(images)
                
                # Save text and images to PDF with Korean font support
                if _save_text_images_to_pdf_korean(full_text.strip(), all_images, output_file):
                    return True, str(output_file)
                else:
                    return False, "Failed to save PDF"
            else:
                return False, "Unable to read HWP document"
                
        except ImportError as e:
            return False, f"{e} not installed"
        except Exception as e:
            error_msg = f"Error during HWP conversion: {str(e)}"
            print(f"[Error] {error_msg}", file=sys.stderr)
            return False, error_msg

    def _convert_docx_to_pdf(self, input_file: Path, output_file: Path) -> Tuple[bool, str]:
        """Convert DOCX/DOC to PDF with multiple fallback methods"""
        # Try LibreOffice first if available
        if self._libreoffice_available:
            print(f"Trying LibreOffice for {input_file.suffix} conversion...", file=sys.stderr)
            success, result = self._convert_with_libreoffice(input_file, output_file)
            if success:
                return success, result
            print(f"LibreOffice conversion failed, falling back to other methods: {result}", file=sys.stderr)
        
        # Fallback to existing methods
        temp_dir = None
        
        try:
            # Create temp directory
            temp_dir = output_file.parent / f"temp_{input_file.stem}"
            temp_dir.mkdir(exist_ok=True)
            
            # Check if it's a .doc file (legacy format)
            if input_file.suffix.lower() == '.doc':
                # Try pure Python OLE parser
                try:
                    success, result = self._convert_doc_with_olefile(input_file, output_file, temp_dir)
                    if success:
                        return success, result
                    print(f"olefile method failed: {result}", file=sys.stderr)
                except Exception as e:
                    print(f"olefile method error: {e}", file=sys.stderr)
                
                # .doc conversion failed
                return False, f"Legacy .doc format could not be processed. Please convert your file to .docx format and try again. You can use Microsoft Word, Google Docs, or LibreOffice to convert .doc to .docx."
            
            # Method 1: Try python-docx for DOCX files
            if input_file.suffix.lower() == '.docx':
                try:
                    success, result = self._convert_docx_with_python_docx(input_file, output_file, temp_dir)
                    if success:
                        return success, result
                    print(f"python-docx method failed: {result}", file=sys.stderr)
                except Exception as e:
                    print(f"python-docx method error: {e}", file=sys.stderr)
            
            # Method 2: Try pypandoc (supports .docx only)
            if input_file.suffix.lower() == '.docx':
                try:
                    success, result = self._convert_with_pypandoc(input_file, output_file, temp_dir)
                    if success:
                        return success, result
                    print(f"pypandoc method failed: {result}", file=sys.stderr)
                except Exception as e:
                    print(f"pypandoc method error: {e}", file=sys.stderr)
            
            # Method 3: Try direct pandoc subprocess (docx only)
            if input_file.suffix.lower() == '.docx':
                try:
                    success, result = self._convert_with_pandoc_subprocess(input_file, output_file, temp_dir)
                    if success:
                        return success, result
                    print(f"pandoc subprocess method failed: {result}", file=sys.stderr)
                except Exception as e:
                    print(f"pandoc subprocess method error: {e}", file=sys.stderr)
            
            # All methods failed
            return False, f"All conversion methods failed for {input_file.suffix} file. Please ensure pandoc is installed and the file is not corrupted."
                
        except Exception as e:
            return False, f"Document conversion error: {str(e)}"
        finally:
            # Cleanup temp directory
            if temp_dir and temp_dir.exists():
                try:
                    import shutil
                    shutil.rmtree(temp_dir)
                except Exception as e:
                    print(f"Cleanup warning: {e}", file=sys.stderr)
    
    def _convert_doc_with_olefile(self, input_file: Path, output_file: Path, temp_dir: Path) -> Tuple[bool, str]:
        """Convert .doc using pure Python olefile parser"""
        try:
            # Try to import olefile
            try:
                import olefile
            except ImportError:
                return False, "olefile not available. Install with: pip install olefile"
            
            # Check if it's a valid OLE file
            if not olefile.isOleFile(str(input_file)):
                return False, "File is not a valid OLE document"
            
            # Open the OLE file
            with olefile.OleFileIO(str(input_file)) as ole:
                # Try to find WordDocument stream (contains the main text)
                if ole._olestream_size('WordDocument') is None:
                    return False, "No WordDocument stream found in .doc file"
                
                # Extract the WordDocument stream
                word_stream = ole._olestream_data('WordDocument')
                
                # Basic text extraction (simplified approach)
                # .doc format is complex, this is a basic attempt
                text_content = ""
                
                # Try to extract readable text from the binary data
                # This is a very basic approach - real .doc parsing is much more complex
                try:
                    # Look for text patterns in the binary data
                    import re
                    
                    # Convert to string and extract readable text
                    # This is a simplified approach and may miss formatting
                    decoded_text = word_stream.decode('latin-1', errors='ignore')
                    
                    # Extract text using regex patterns
                    # Look for sequences of printable characters
                    text_matches = re.findall(r'[a-zA-Z가-힣\s.,!?;:()"\'-]{10,}', decoded_text)
                    
                    if text_matches:
                        text_content = '\n'.join(text_matches)
                        # Clean up the text
                        text_content = re.sub(r'\s+', ' ', text_content)
                        text_content = text_content.strip()
                    
                    if not text_content:
                        # Try another approach with different encoding
                        try:
                            decoded_text = word_stream.decode('utf-8', errors='ignore')
                            text_matches = re.findall(r'[a-zA-Z가-힣\s.,!?;:()"\'-]{10,}', decoded_text)
                            if text_matches:
                                text_content = '\n'.join(text_matches)
                                text_content = re.sub(r'\s+', ' ', text_content)
                                text_content = text_content.strip()
                        except:
                            pass
                    
                except Exception as e:
                    print(f"Text extraction error: {e}", file=sys.stderr)
                
                if not text_content.strip():
                    return False, "No readable text could be extracted from .doc file"
                
                # Convert extracted text to PDF
                if _save_text_images_to_pdf_korean(text_content, [], output_file):
                    return True, str(output_file)
                else:
                    return False, "Failed to save PDF from extracted text"
                    
        except Exception as e:
            return False, f"olefile conversion error: {str(e)}"
    
    def _convert_docx_with_python_docx(self, input_file: Path, output_file: Path, temp_dir: Path) -> Tuple[bool, str]:
        """Convert DOCX using python-docx library"""
        try:
            from docx import Document
            from docx.shared import Inches
            import io
            
            # Load document
            doc = Document(str(input_file))
            
            # Extract text
            full_text = ""
            image_files = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    full_text += para.text + "\n\n"
            
            # Extract images from document
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    try:
                        image_data = rel.target_part.blob
                        
                        # Determine image extension
                        if image_data[:8] == b"\x89PNG\r\n\x1a\n":
                            ext = ".png"
                        elif image_data[:2] == b"\xff\xd8":
                            ext = ".jpg"
                        elif image_data[:6] == b"GIF87a" or image_data[:6] == b"GIF89a":
                            ext = ".gif"
                        else:
                            ext = ".png"  # Default
                        
                        img_path = temp_dir / f"extracted_img_{len(image_files)}{ext}"
                        with open(img_path, 'wb') as f:
                            f.write(image_data)
                        image_files.append(str(img_path))
                    except Exception as e:
                        print(f"Image extraction warning: {e}", file=sys.stderr)
            
            # Extract tables
            for table in doc.tables:
                full_text += "\n[Table]\n"
                for row in table.rows:
                    row_text = " | ".join([cell.text for cell in row.cells])
                    full_text += row_text + "\n"
                full_text += "\n"
            
            if not full_text.strip():
                full_text = "No readable text content found in document."
            
            # Use Korean PDF function
            if _save_text_images_to_pdf_korean(full_text, image_files, output_file):
                return True, str(output_file)
            else:
                return False, "Failed to save PDF"
                
        except ImportError:
            return False, "python-docx not available. Install with: pip install python-docx"
        except Exception as e:
            return False, f"python-docx conversion failed: {str(e)}"
    
    def _convert_with_pypandoc(self, input_file: Path, output_file: Path, temp_dir: Path) -> Tuple[bool, str]:
        """Convert using pypandoc library"""
        try:
            # Ensure pandoc is available
            pypandoc.ensure_pandoc_installed()
            
            # Extract text using pypandoc
            full_text = pypandoc.convert_file(str(input_file), 'plain', encoding='utf-8')
            
            if not full_text.strip():
                return False, "No text content extracted"
            
            # Extract images using pypandoc
            image_files = []
            try:
                # Get pandoc path from pypandoc
                pandoc_path = pypandoc.get_pandoc_path()
                
                import subprocess
                result = subprocess.run([
                    pandoc_path, str(input_file),
                    '-t', 'markdown',
                    '-o', str(temp_dir / 'dummy.md'),
                    '--extract-media', str(temp_dir)
                ], capture_output=True, text=True, timeout=30)
                
                # Find extracted images
                media_dir = temp_dir / "media"
                if media_dir.exists():
                    for img_file in media_dir.rglob('*'):
                        if img_file.suffix.lower() in ['.png', '.jpg', '.jpeg', '.jfif', '.gif', '.bmp', '.tiff']:
                            image_files.append(str(img_file))
            except Exception as e:
                print(f"Image extraction warning: {e}", file=sys.stderr)
            
            # Use Korean PDF function
            if _save_text_images_to_pdf_korean(full_text, image_files, output_file):
                return True, str(output_file)
            else:
                return False, "Failed to save PDF"
                
        except Exception as e:
            return False, f"pypandoc conversion failed: {str(e)}"
    
    def _convert_with_pandoc_subprocess(self, input_file: Path, output_file: Path, temp_dir: Path) -> Tuple[bool, str]:
        """Convert using pandoc subprocess directly"""
        try:
            import subprocess
            import shutil
            
            # Check if pandoc is available
            if not shutil.which('pandoc'):
                return False, "pandoc not found in PATH. Please install pandoc."
            
            # Extract text
            result = subprocess.run([
                'pandoc', str(input_file), 
                '-t', 'plain',
                '--wrap=none'
            ], capture_output=True, text=True, timeout=60)
            
            if result.returncode != 0:
                error_msg = result.stderr.strip() if result.stderr else "pandoc text extraction failed"
                return False, f"pandoc error: {error_msg}"
            
            full_text = result.stdout
            if not full_text.strip():
                return False, "No text content extracted by pandoc"
            
            # Extract images
            image_files = []
            try:
                result = subprocess.run([
                    'pandoc', str(input_file),
                    '-t', 'markdown',
                    '-o', str(temp_dir / 'dummy.md'),
                    '--extract-media', str(temp_dir)
                ], capture_output=True, text=True, timeout=60)
                
                # Find extracted images
                media_dir = temp_dir / "media"
                if media_dir.exists():
                    for img_file in media_dir.rglob('*'):
                        if img_file.suffix.lower() in ['.png', '.jpg', '.jpeg', '.jfif', '.gif', '.bmp', '.tiff']:
                            image_files.append(str(img_file))
            except Exception as e:
                print(f"Image extraction warning: {e}", file=sys.stderr)
            
            # Use Korean PDF function
            if _save_text_images_to_pdf_korean(full_text, image_files, output_file):
                return True, str(output_file)
            else:
                return False, "Failed to save PDF"
                
        except subprocess.TimeoutExpired:
            return False, "pandoc conversion timed out (file may be too large or corrupted)"
        except FileNotFoundError:
            return False, "pandoc executable not found. Please install pandoc."
        except Exception as e:
            return False, f"pandoc subprocess failed: {str(e)}"

    def _convert_excel_to_pdf(self, input_file: Path, output_file: Path) -> Tuple[bool, str]:
        """Convert Excel to PDF using openpyxl"""
        # Try LibreOffice first if available
        if self._libreoffice_available:
            print(f"Trying LibreOffice for {input_file.suffix} conversion...", file=sys.stderr)
            success, result = self._convert_with_libreoffice(input_file, output_file)
            if success:
                return success, result
            print(f"LibreOffice conversion failed, falling back to openpyxl: {result}", file=sys.stderr)
        
        # Fallback to openpyxl
        try:
            import openpyxl
            
            # Load Excel file
            wb = openpyxl.load_workbook(input_file, data_only=True)
            
            # Extract all text content
            full_text = ""
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                full_text += f"=== Sheet: {sheet_name} ===\n\n"
                
                # Extract data from sheet
                for row in sheet.iter_rows(values_only=True):
                    row_text = "\t".join([str(cell) if cell is not None else '' for cell in row])
                    if row_text.strip():
                        full_text += row_text + "\n"
                
                full_text += "\n\n"
            
            images = []

            if _save_text_images_to_pdf_korean(full_text.strip(), images, output_file):
                return True, str(output_file)
            else:
                return False, "Failed to save PDF"
                
        except ImportError:
            return False, "Please install openpyxl: pip install openpyxl"
        except Exception as e:
            return False, f"Excel conversion error: {str(e)}"

    def _convert_ppt_to_pdf(self, input_file: Path, output_file: Path) -> Tuple[bool, str]:
        """Convert PowerPoint to PDF using python-pptx"""
        # Try LibreOffice first if available
        if self._libreoffice_available:
            print(f"Trying LibreOffice for {input_file.suffix} conversion...", file=sys.stderr)
            success, result = self._convert_with_libreoffice(input_file, output_file)
            if success:
                return success, result
            print(f"LibreOffice conversion failed, falling back to python-pptx: {result}", file=sys.stderr)
        
        # Fallback to python-pptx
        try:
            from pptx import Presentation
            from PIL import Image
            import io
            
            # Load PowerPoint file
            prs = Presentation(str(input_file))
            
            # Extract text and images
            full_text = ""
            image_paths = []
            temp_dir = output_file.parent / f"temp_{input_file.stem}"
            temp_dir.mkdir(exist_ok=True)
            
            for slide_num, slide in enumerate(prs.slides):
                full_text += f"=== Slide {slide_num + 1} ===\n\n"
                
                # Extract text from slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        full_text += shape.text.strip() + "\n\n"
                    
                    # Extract images
                    if shape.shape_type == 13:  # Picture
                        try:
                            image = shape.image
                            image_bytes = image.blob
                            img = Image.open(io.BytesIO(image_bytes))
                            
                            # Save image
                            img_path = temp_dir / f"slide_{slide_num}_img_{shape.shape_id}.png"
                            img.save(img_path)
                            image_paths.append(str(img_path))
                        except:
                            pass
                
                # Extract tables
                for shape in slide.shapes:
                    if shape.has_table:
                        full_text += "[Table]\n"
                        table = shape.table
                        for row in table.rows:
                            row_text = " | ".join([cell.text for cell in row.cells])
                            full_text += row_text + "\n"
                        full_text += "\n"
                
                full_text += "\n"
            
            if _save_text_images_to_pdf_korean(full_text.strip(), image_paths, output_file):
                # Cleanup temp directory
                import shutil
                shutil.rmtree(temp_dir)
                return True, str(output_file)
            else:
                return False, "Failed to save PDF"
                
        except ImportError:
            return False, "Please install python-pptx and pillow: pip install python-pptx pillow"
        except Exception as e:
            return False, f"PowerPoint conversion error: {str(e)}"
            
    def _convert_text_to_pdf(self, input_file: Path, output_file: Path) -> Tuple[bool, str]:
        """Convert plain text to PDF"""
        try:
            # Read text file
            with open(input_file, 'r', encoding='utf-8') as f:
                text = f.read()

            try:
                doc = SimpleDocTemplate(str(output_file), pagesize=A4)
                styles = getSampleStyleSheet()
                story = []
                
                # Title
                title_style = ParagraphStyle(
                    'CustomTitle',
                    parent=styles['Heading1'],
                    fontSize=24,
                    spaceAfter=30
                )
                story.append(Paragraph(input_file.name, title_style))
                story.append(Spacer(1, 0.2*inch))
                
                # Content
                text_style = ParagraphStyle(
                    'CustomText',
                    parent=styles['Normal'],
                    fontSize=11,
                    leading=14
                )
                
                # Split text into paragraphs
                for paragraph in text.split('\n\n'):
                    if paragraph.strip():
                        # Escape special characters for reportlab
                        safe_text = paragraph.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                        story.append(Paragraph(safe_text, text_style))
                        story.append(Spacer(1, 0.1*inch))
                
                doc.build(story)
                return True, str(output_file)
                
            except Exception as e:
                print(f"Reportlab conversion failed: {e}", file=sys.stderr)
            
            # Second try: Direct PDF creation using pure Python
            try:
                # Simple PDF creation without external dependencies
                pdf_content = self._create_simple_pdf(text, input_file.name)
                with open(output_file, 'wb') as f:
                    f.write(pdf_content)
                return True, str(output_file)
            except Exception as e:
                print(f"Simple PDF creation failed: {e}", file=sys.stderr)
            
            # Third try: Create HTML and convert
            html_content = f"""
            <html>
            <head>
                <meta charset="utf-8">
                <style>
                    body {{
                        font-family: Arial, sans-serif;
                        margin: 40px;
                        line-height: 1.6;
                        white-space: pre-wrap;
                    }}
                    h1 {{
                        color: #333;
                        border-bottom: 2px solid #333;
                        padding-bottom: 10px;
                    }}
                    pre {{
                        background: #f4f4f4;
                        padding: 15px;
                        border-radius: 5px;
                        overflow-x: auto;
                    }}
                </style>
            </head>
            <body>
                <h1>{input_file.name}</h1>
                <pre>{text}</pre>
            </body>
            </html>
            """
            
            # Try various HTML to PDF converters
            converters_tried = []
            # Try pdfkit
            try:
                pdfkit.from_string(html_content, str(output_file))
                return True, str(output_file)
            except Exception as e:
                converters_tried.append("pdfkit")
            
            try:
                temp_html = output_file.with_suffix('.html')
                with open(temp_html, 'w', encoding='utf-8') as f:
                    f.write(html_content)
                
                pypandoc.convert_file(str(temp_html), 'pdf', outputfile=str(output_file))
                temp_html.unlink()
                return True, str(output_file)
            except Exception as e:
                converters_tried.append("pandoc")
                    
        except Exception as e:
            print(f"Text to PDF conversion error: {e}", file=sys.stderr)
        
        # Provide helpful error message
        error_msg = "Text conversion failed. "
        if converters_tried:
            error_msg += f"Tried: {', '.join(converters_tried)}. "
        else:
            error_msg += "No PDF converters available. Install: pip install reportlab"
        
        return False, error_msg

    def _create_simple_pdf(self, text: str, title: str) -> bytes:
        """Create a very simple PDF without external dependencies"""
        # This is a minimal PDF creator - for production use reportlab is recommended
        lines = text.split('\n')
        
        # Basic PDF structure
        pdf = b"%PDF-1.4\n"
        
        # Catalog and Pages
        pdf += b"1 0 obj\n<< /Type /Catalog /Pages 2 0 R >>\nendobj\n"
        pdf += b"2 0 obj\n<< /Type /Pages /Kids [3 0 R] /Count 1 >>\nendobj\n"
        
        # Page
        pdf += b"3 0 obj\n<< /Type /Page /Parent 2 0 R /Resources 4 0 R /MediaBox [0 0 612 792] /Contents 5 0 R >>\nendobj\n"
        
        # Resources
        pdf += b"4 0 obj\n<< /Font << /F1 << /Type /Font /Subtype /Type1 /BaseFont /Helvetica >> >> >>\nendobj\n"
        
        # Content stream
        content = f"BT /F1 12 Tf 50 750 Td ({title}) Tj ET\n"
        y_pos = 720
        for line in lines[:50]:  # Limit to first 50 lines for simplicity
            if line.strip():
                # Escape special characters
                safe_line = line.replace('\\', '\\\\').replace('(', '\\(').replace(')', '\\)')
                content += f"BT /F1 10 Tf 50 {y_pos} Td ({safe_line[:80]}) Tj ET\n"
                y_pos -= 15
                if y_pos < 50:
                    break
        
        content_bytes = content.encode('latin-1', errors='replace')
        pdf += f"5 0 obj\n<< /Length {len(content_bytes)} >>\nstream\n".encode()
        pdf += content_bytes
        pdf += b"\nendstream\nendobj\n"
        
        # xref table
        xref_pos = len(pdf)
        pdf += b"xref\n0 6\n"
        pdf += b"0000000000 65535 f \n"
        pdf += b"0000000009 00000 n \n"
        pdf += b"0000000058 00000 n \n"
        pdf += b"0000000115 00000 n \n"
        pdf += b"0000000229 00000 n \n"
        pdf += b"0000000328 00000 n \n"
        
        # Trailer
        pdf += b"trailer\n<< /Size 6 /Root 1 0 R >>\n"
        pdf += f"startxref\n{xref_pos}\n".encode()
        pdf += b"%%EOF"
        
        return pdf
    
    def _convert_image_to_pdf(self, input_file: Path, output_file: Path) -> Tuple[bool, str]:
        """Convert image to PDF with high quality preservation"""
        # Try LibreOffice first if available
        if self._libreoffice_available:
            print(f"Trying LibreOffice for image conversion...", file=sys.stderr)
            success, result = self._convert_with_libreoffice(input_file, output_file)
            if success:
                return success, result
            print(f"LibreOffice conversion failed, falling back to PIL: {result}", file=sys.stderr)
        
        # Fallback to PIL-based conversion
        try:
            # Open image
            img = Image.open(input_file)
            
            # Get original DPI if available, default to 300 for high quality
            original_dpi = img.info.get('dpi', (300, 300))
            if isinstance(original_dpi, tuple):
                dpi = max(original_dpi[0], 300)  # Use at least 300 DPI
            else:
                dpi = max(float(original_dpi), 300)
            
            # Log image info
            print(f"Converting image: {input_file.name}", file=sys.stderr)
            print(f"  Original size: {img.size}", file=sys.stderr)
            print(f"  Original DPI: {original_dpi}", file=sys.stderr)
            print(f"  Using DPI: {dpi}", file=sys.stderr)
            
            # Convert to RGB if necessary (preserving quality)
            if img.mode in ('RGBA', 'LA'):
                # Create white background
                background = Image.new('RGB', img.size, (255, 255, 255))
                if img.mode == 'RGBA':
                    background.paste(img, mask=img.split()[-1])
                else:
                    background.paste(img)
                img = background
            elif img.mode == 'P':
                # Convert palette images to RGB for better quality
                img = img.convert('RGB')
            elif img.mode not in ('RGB', 'L', 'CMYK'):
                img = img.convert('RGB')
            
            # Save as PDF with high quality
            # Use high quality settings
            img.save(
                str(output_file), 
                'PDF', 
                resolution=dpi,
                quality=95,  # High JPEG quality for embedded images
                optimize=True,  # Optimize file size without quality loss
                dpi=(dpi, dpi)  # Explicitly set DPI metadata
            )
            
            print(f"  Converted to PDF with {dpi} DPI", file=sys.stderr)
            return True, str(output_file)
            
        except Exception as e:
            return False, f"Image conversion failed: {e}"
    
    def _convert_with_pandoc_simple(self, input_file: Path, output_file: Path) -> Tuple[bool, str]:
        """Convert documents using Pandoc without LaTeX"""
        temp_dir = None
        try:
            import subprocess
            import shutil
            
            # Check if pandoc is available
            if not shutil.which('pandoc'):
                return False, "pandoc not found in PATH. Please install pandoc."
            
            # Create temp directory for extracted images
            temp_dir = output_file.parent / f"temp_{input_file.stem}"
            temp_dir.mkdir(exist_ok=True)
            
            # Extract text with better error handling
            result = subprocess.run([
                'pandoc', str(input_file), 
                '-t', 'plain',
                '--wrap=none'
            ], capture_output=True, text=True, timeout=60)
            
            if result.returncode != 0:
                error_msg = result.stderr.strip() if result.stderr else "pandoc text extraction failed"
                return False, f"pandoc error: {error_msg}"
            
            full_text = result.stdout
            if not full_text.strip():
                return False, "No text content extracted by pandoc"
            
            # Extract images with timeout
            image_files = []
            try:
                result = subprocess.run([
                    'pandoc', str(input_file),
                    '-t', 'markdown',
                    '-o', str(temp_dir / 'dummy.md'),
                    '--extract-media', str(temp_dir)
                ], capture_output=True, text=True, timeout=60)
                
                # Find extracted images
                media_dir = temp_dir / "media"
                if media_dir.exists():
                    for img_file in media_dir.rglob('*'):
                        if img_file.suffix.lower() in ['.png', '.jpg', '.jpeg', '.jfif', '.gif', '.bmp', '.tiff']:
                            image_files.append(str(img_file))
            except subprocess.TimeoutExpired:
                print(f"Image extraction timed out, continuing without images", file=sys.stderr)
            except Exception as e:
                print(f"Image extraction warning: {e}", file=sys.stderr)
            
            # Use Korean-supporting PDF function
            if _save_text_images_to_pdf_korean(full_text, image_files, output_file):
                return True, str(output_file)
            else:
                return False, "Failed to save PDF"
                
        except subprocess.TimeoutExpired:
            return False, "pandoc conversion timed out (file may be too large or corrupted)"
        except FileNotFoundError:
            return False, "pandoc executable not found. Please install pandoc."
        except Exception as e:
            return False, f"Conversion error: {str(e)}"
        finally:
            # Cleanup temp directory
            if temp_dir and temp_dir.exists():
                try:
                    import shutil
                    shutil.rmtree(temp_dir)
                except Exception as e:
                    print(f"Cleanup warning: {e}", file=sys.stderr)

    def _convert_audio_to_pdf(self, input_file: Path, output_file: Path) -> Tuple[bool, str]:
        """Convert audio file to PDF with transcription"""
        if not MULTIMEDIA_SUPPORT:
            return False, "Audio conversion requires whisper. Install with: pip install openai-whisper"
        
        try:
            print(f"Converting audio file to PDF with transcription...", file=sys.stderr)
            
            # Process audio file
            txt_path = process_audio_file(str(input_file), str(output_file.parent))
            
            # Convert the transcription text file to PDF
            return self._convert_text_to_pdf(Path(txt_path), output_file)
            
        except Exception as e:
            return False, f"Audio conversion failed: {str(e)}"
    
    def _convert_video_to_pdf(self, input_file: Path, output_file: Path) -> Tuple[bool, str]:
        """Convert video file to PDF with frames and transcription"""
        if not MULTIMEDIA_SUPPORT:
            return False, "Video conversion requires whisper and ffmpeg. Install with: pip install openai-whisper && install ffmpeg"
        
        try:
            print(f"Converting video file to PDF with frames and transcription...", file=sys.stderr)
            print(f"This may take several minutes for large videos...", file=sys.stderr)
            
            # Process video file
            pdf_path = process_video_file(str(input_file), str(output_file.parent))
            
            # If the output path is different from what we got, move the file
            if Path(pdf_path) != output_file:
                import shutil
                shutil.move(pdf_path, str(output_file))
                return True, str(output_file)
            else:
                return True, pdf_path
            
        except Exception as e:
            return False, f"Video conversion failed: {str(e)}"

def convert_file_to_pdf(input_path: str, output_dir: Optional[str] = None) -> Tuple[bool, str]:
    """
    Convenience function to convert a file to PDF
    
    Args:
        input_path: Path to input file
        output_dir: Optional output directory
        
    Returns:
        Tuple of (success: bool, output_path_or_error: str)
    """
    converter = FileConverter(Path(output_dir) if output_dir else None)
    return converter.convert_to_pdf(input_path)


if __name__ == "__main__":
    # Test conversion
    import argparse
    
    parser = argparse.ArgumentParser(description="Convert files to PDF")
    parser.add_argument("input_file", help="Input file path")
    parser.add_argument("-o", "--output", help="Output PDF path")
    parser.add_argument("-d", "--output-dir", help="Output directory")
    
    args = parser.parse_args()
    
    converter = FileConverter(Path(args.output_dir) if args.output_dir else None)
    success, result = converter.convert_to_pdf(args.input_file, args.output)
    
    if success:
        print(f"✅ Converted successfully: {result}", file=sys.stderr)
    else:
        print(f"❌ Conversion failed: {result}", file=sys.stderr)
        sys.exit(1)
